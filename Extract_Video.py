#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
🎬 올인원 비디오 처리기 v4.0 - 완전 리팩토링 버전D
매니저 패턴, 한글 지원, 개선된 에러 핸들링

주요 개선사항:
- 매니저 패턴으로 코드 구조 개선
- 한글 파일명/경로 완전 지원
- 싱글턴 패턴으로 리소스 관리 최적화
- 샘플링 간격 조절 기능 추가
- 강화된 에러 핸들링 및 복구 로직
"""

import os
import sys
import numpy as np
import logging
import tempfile
import re
import json
import subprocess
import threading
import queue
import argparse
import math
import time
import gc
import atexit
import shutil
import uuid
from pathlib import Path
from datetime import datetime, timedelta
from typing import Tuple, List, Optional, Dict, Any, Union
import warnings
from contextlib import contextmanager
from collections import deque
from dataclasses import dataclass, field
from enum import Enum


# ============================================================================
# 라이브러리 Import 및 검증
# ============================================================================

def safe_import(module_name, package_name=None, pip_name=None):
    """안전한 라이브러리 import"""
    try:
        if package_name:
            module = __import__(module_name, fromlist=[package_name])
            return getattr(module, package_name), True
        else:
            return __import__(module_name), True
    except ImportError:
        pip_name = pip_name or module_name
        warnings.warn(f"{module_name}을 import할 수 없습니다. 설치: pip install {pip_name}")
        return None, False


# 필수 라이브러리
try:
    import cv2

    CV2_AVAILABLE = True
except ImportError:
    print("Error: OpenCV가 설치되지 않았습니다. (pip install opencv-python)")
    sys.exit(1)

try:
    from PIL import Image

    PIL_AVAILABLE = True
except ImportError:
    print("Error: Pillow가 설치되지 않았습니다. (pip install Pillow)")
    sys.exit(1)

# 선택적 라이브러리
GUI_AVAILABLE = False
try:
    import tkinter as tk
    from tkinter import ttk, filedialog, messagebox, scrolledtext

    GUI_AVAILABLE = True
except ImportError:
    warnings.warn("tkinter를 사용할 수 없습니다. CLI 모드만 사용 가능합니다.")

# 기타 라이브러리들
SSIM_AVAILABLE = False
try:
    from skimage.metrics import structural_similarity as ssim

    SSIM_AVAILABLE = True
except ImportError:
    warnings.warn("scikit-image가 없어 기본 유사도 계산을 사용합니다.")

pptx_module, PPTX_AVAILABLE = safe_import('pptx', pip_name='python-pptx')
if PPTX_AVAILABLE:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

whisper_module, WHISPER_AVAILABLE = safe_import('whisper', pip_name='openai-whisper')
if WHISPER_AVAILABLE:
    import whisper

moviepy_module, MOVIEPY_AVAILABLE = safe_import('moviepy.editor', package_name='VideoFileClip', pip_name='moviepy')
if MOVIEPY_AVAILABLE:
    from moviepy.editor import VideoFileClip

docx_module, DOCX_AVAILABLE = safe_import('docx', pip_name='python-docx')
if DOCX_AVAILABLE:
    from docx import Document
    from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

# SpellChecker 및 OCRManager import
try:
    from spell_checker import SpellChecker
    from ocr_manager import OCRManager
    SPELL_CHECKER_AVAILABLE = True
    OCR_AVAILABLE = True
except ImportError as e:
    print(f"Warning: SpellChecker/OCRManager를 import할 수 없습니다: {e}")
    SPELL_CHECKER_AVAILABLE = False
    OCR_AVAILABLE = False


# ============================================================================
# 열거형 및 데이터 클래스
# ============================================================================

class ProcessingStatus(Enum):
    """처리 상태"""
    IDLE = "idle"
    PROCESSING = "processing"
    COMPLETED = "completed"
    ERROR = "error"
    CANCELLED = "cancelled"


class LogLevel(Enum):
    """로그 레벨"""
    DEBUG = logging.DEBUG
    INFO = logging.INFO
    WARNING = logging.WARNING
    ERROR = logging.ERROR
    CRITICAL = logging.CRITICAL


@dataclass
class ProcessingConfig:
    """처리 설정"""
    similarity_threshold: float = 0.95
    adaptive_threshold: bool = True
    extract_frames: bool = True
    create_ppt: bool = True
    extract_audio: bool = True
    create_srt: bool = True
    create_word: bool = True
    ocr_spell_check: bool = False  # OCR 맞춤법 검수 (기본 비활성화 - 시간 소요)
    grid_rows: int = 3
    grid_cols: int = 3
    whisper_model: str = "base"
    frame_sampling_interval: float = 1.0  # 새로운 설정: 샘플링 간격 (초)
    min_frame_interval: float = 0.5  # 최소 프레임 간격
    max_frames_per_video: int = 200  # 비디오당 최대 프레임 수


@dataclass
class ProcessingResult:
    """처리 결과"""
    video_path: str
    output_dir: Optional[str] = None
    frames: List[str] = field(default_factory=list)
    frame_count: int = 0
    ppt_path: Optional[str] = None
    audio_path: Optional[str] = None
    srt_path: Optional[str] = None
    word_path: Optional[str] = None
    transcription: Optional[Dict] = None
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    processing_time: float = 0.0
    status: ProcessingStatus = ProcessingStatus.IDLE


# ============================================================================
# 싱글턴 메타클래스
# ============================================================================

class SingletonMeta(type):
    """싱글턴 메타클래스"""
    _instances = {}
    _lock = threading.Lock()

    def __call__(cls, *args, **kwargs):
        if cls not in cls._instances:
            with cls._lock:
                if cls not in cls._instances:
                    cls._instances[cls] = super().__call__(*args, **kwargs)
        return cls._instances[cls]


# ============================================================================
# 에러 매니저
# ============================================================================

class ErrorManager(metaclass=SingletonMeta):
    """에러 핸들링 및 로깅 관리"""

    def __init__(self):
        self.error_counts = {}
        self.max_retries = 3
        self.retry_delays = [1, 2, 4]  # 지수백오프
        self.setup_logging()

    def setup_logging(self):
        """로깅 설정"""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            handlers=[
                logging.StreamHandler(sys.stdout),
                logging.FileHandler('video_processor_v4.log', encoding='utf-8')
            ]
        )
        self.logger = logging.getLogger('VideoProcessor')

    def log(self, level: LogLevel, message: str, exception: Exception = None):
        """로그 메시지 기록"""
        if exception:
            message = f"{message}: {str(exception)}"

        self.logger.log(level.value, message)

    def handle_error(self, error_key: str, exception: Exception,
                     context: str = "") -> bool:
        """에러 처리 및 재시도 로직"""
        full_key = f"{error_key}_{context}" if context else error_key

        if full_key not in self.error_counts:
            self.error_counts[full_key] = 0

        self.error_counts[full_key] += 1
        retry_count = self.error_counts[full_key]

        self.log(LogLevel.ERROR,
                 f"오류 발생 ({retry_count}/{self.max_retries}): {error_key}",
                 exception)

        if retry_count < self.max_retries:
            delay = self.retry_delays[min(retry_count - 1, len(self.retry_delays) - 1)]
            self.log(LogLevel.INFO, f"{delay}초 후 재시도...")
            time.sleep(delay)
            return True  # 재시도 가능
        else:
            self.log(LogLevel.ERROR, f"최대 재시도 횟수 초과: {error_key}")
            return False  # 재시도 불가

    def reset_error_count(self, error_key: str, context: str = ""):
        """에러 카운트 리셋"""
        full_key = f"{error_key}_{context}" if context else error_key
        if full_key in self.error_counts:
            del self.error_counts[full_key]


# ============================================================================
# 설정 매니저
# ============================================================================

class ConfigManager(metaclass=SingletonMeta):
    """설정 관리"""

    def __init__(self):
        self.config = ProcessingConfig()
        self.config_file = Path.home() / '.video_processor_config.json'
        self.load_config()

    def load_config(self):
        """설정 파일 로드"""
        try:
            if self.config_file.exists():
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    for key, value in data.items():
                        if hasattr(self.config, key):
                            setattr(self.config, key, value)
                ErrorManager().log(LogLevel.INFO, "설정 파일 로드 완료")
        except Exception as e:
            ErrorManager().log(LogLevel.WARNING, "설정 파일 로드 실패", e)

    def save_config(self):
        """설정 파일 저장"""
        try:
            config_data = {
                key: getattr(self.config, key)
                for key in dir(self.config)
                if not key.startswith('_')
            }

            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(config_data, f, indent=2, ensure_ascii=False)

            ErrorManager().log(LogLevel.INFO, "설정 파일 저장 완료")
        except Exception as e:
            ErrorManager().log(LogLevel.WARNING, "설정 파일 저장 실패", e)

    def get_config(self) -> ProcessingConfig:
        """현재 설정 반환"""
        return self.config

    def update_config(self, **kwargs):
        """설정 업데이트"""
        for key, value in kwargs.items():
            if hasattr(self.config, key):
                setattr(self.config, key, value)
        self.save_config()


# ============================================================================
# 경로 매니저 - 한글 지원 강화
# ============================================================================

class PathManager(metaclass=SingletonMeta):
    """경로 및 파일명 처리 - 한글 지원"""

    def __init__(self):
        self.temp_files = []
        self.temp_dirs = []
        self.encoding = 'utf-8'
        atexit.register(self.cleanup_all)

    def normalize_path(self, path: Union[str, Path]) -> Path:
        """경로 정규화"""
        try:
            return Path(path).resolve()
        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, f"경로 정규화 실패: {path}", e)
            return Path(str(path))

    def safe_filename(self, filename: str, max_length: int = 200) -> str:
        """안전한 파일명 생성 - 한글 보존"""
        if not filename:
            return "unnamed"

        # 금지된 문자만 제거 (한글은 보존)
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        filename = filename.strip('. ')

        # 길이 제한
        if len(filename.encode('utf-8')) > max_length:
            # UTF-8 바이트 길이 기준으로 자르기
            encoded = filename.encode('utf-8')
            truncated = encoded[:max_length]
            # 깨진 문자 방지
            try:
                filename = truncated.decode('utf-8')
            except UnicodeDecodeError:
                # 마지막 불완전한 문자 제거
                while len(truncated) > 0:
                    try:
                        filename = truncated.decode('utf-8')
                        break
                    except UnicodeDecodeError:
                        truncated = truncated[:-1]
                else:
                    filename = "unnamed"

        # Windows 예약어 확인
        name_only = filename.split('.')[0].upper()
        reserved = ['CON', 'PRN', 'AUX', 'NUL'] + \
                   [f'COM{i}' for i in range(1, 10)] + \
                   [f'LPT{i}' for i in range(1, 10)]

        if name_only in reserved:
            filename = f"_{filename}"

        return filename or "unnamed"

    def create_temp_file(self, suffix='.tmp', prefix='video_proc_',
                         content: bytes = None) -> Optional[Path]:
        """임시 파일 생성 - 한글 경로 처리"""
        try:
            # 시스템 임시 디렉토리 사용
            temp_dir = Path(tempfile.gettempdir())

            # 유니크한 파일명 생성 (영문으로)
            unique_name = f"{prefix}{uuid.uuid4().hex[:8]}{suffix}"
            temp_path = temp_dir / unique_name

            # 파일 생성
            if content is not None:
                temp_path.write_bytes(content)
            else:
                temp_path.touch()

            self.temp_files.append(temp_path)
            return temp_path

        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, "임시 파일 생성 실패", e)
            return None

    def create_temp_dir(self, prefix='video_proc_') -> Optional[Path]:
        """임시 디렉토리 생성"""
        try:
            temp_dir = Path(tempfile.mkdtemp(prefix=prefix))
            self.temp_dirs.append(temp_dir)
            return temp_dir
        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, "임시 디렉토리 생성 실패", e)
            return None

    def ensure_directory(self, path: Union[str, Path]) -> bool:
        """디렉토리 생성 확인"""
        try:
            path = self.normalize_path(path)
            path.mkdir(parents=True, exist_ok=True)
            return True
        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, f"디렉토리 생성 실패: {path}", e)
            return False

    def copy_with_korean_support(self, src: Path, dst: Path) -> bool:
        """한글 경로 지원하는 파일 복사"""
        try:
            # 먼저 임시 파일로 복사한 후 최종 위치로 이동
            if src.exists():
                dst.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(str(src), str(dst))
                return True
            return False
        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, f"파일 복사 실패: {src} -> {dst}", e)
            return False

    def read_image_korean(self, image_path: Union[str, Path]) -> Optional[np.ndarray]:
        """한글 경로 이미지 읽기"""
        try:
            image_path = self.normalize_path(image_path)

            # OpenCV 한글 경로 문제 우회: 바이트로 읽어서 디코딩
            with open(image_path, 'rb') as f:
                image_data = f.read()

            # NumPy array로 변환 후 OpenCV로 디코딩
            nparr = np.frombuffer(image_data, np.uint8)
            image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

            return image

        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, f"이미지 읽기 실패: {image_path}", e)
            return None

    def write_image_korean(self, image: np.ndarray, image_path: Union[str, Path],
                           quality: int = 90) -> bool:
        """한글 경로 이미지 저장"""
        try:
            image_path = self.normalize_path(image_path)
            image_path.parent.mkdir(parents=True, exist_ok=True)

            # OpenCV 한글 경로 문제 우회: 메모리에서 인코딩 후 저장
            ext = image_path.suffix.lower()
            if ext in ['.jpg', '.jpeg']:
                encode_param = [cv2.IMWRITE_JPEG_QUALITY, quality]
            elif ext == '.png':
                encode_param = [cv2.IMWRITE_PNG_COMPRESSION, 9]
            else:
                encode_param = []

            success, encoded_img = cv2.imencode(ext, image, encode_param)
            if success:
                with open(image_path, 'wb') as f:
                    f.write(encoded_img.tobytes())
                return True

            return False

        except Exception as e:
            ErrorManager().log(LogLevel.ERROR, f"이미지 저장 실패: {image_path}", e)
            return False

    def cleanup_all(self):
        """모든 임시 파일/디렉토리 정리"""
        for temp_path in self.temp_files:
            try:
                if temp_path.exists():
                    temp_path.unlink()
            except Exception as e:
                ErrorManager().log(LogLevel.WARNING, f"임시 파일 삭제 실패: {temp_path}", e)

        for temp_dir in self.temp_dirs:
            try:
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
            except Exception as e:
                ErrorManager().log(LogLevel.WARNING, f"임시 디렉토리 삭제 실패: {temp_dir}", e)

        self.temp_files.clear()
        self.temp_dirs.clear()


# ============================================================================
# 비디오 매니저
# ============================================================================

class VideoManager:
    """비디오 처리 관리"""

    def __init__(self):
        self.path_manager = PathManager()
        self.error_manager = ErrorManager()
        self.config_manager = ConfigManager()

    @contextmanager
    def safe_video_capture(self, video_path: Union[str, Path]):
        """안전한 VideoCapture 컨텍스트 매니저 - 한글 지원"""
        cap = None
        try:
            video_path = self.path_manager.normalize_path(video_path)

            # 한글 경로 문제 해결을 위해 임시 복사 시도
            if not video_path.exists():
                raise FileNotFoundError(f"비디오 파일이 없습니다: {video_path}")

            # OpenCV가 한글 경로를 처리할 수 있는지 테스트
            cap = cv2.VideoCapture(str(video_path))

            if not cap.isOpened():
                # 한글 경로 문제일 가능성 - 임시 파일로 복사
                temp_video = self.path_manager.create_temp_file(
                    suffix=video_path.suffix, prefix='temp_video_'
                )
                if temp_video and self.path_manager.copy_with_korean_support(video_path, temp_video):
                    cap.release()
                    cap = cv2.VideoCapture(str(temp_video))

                if not cap.isOpened():
                    raise ValueError(f"비디오 파일을 열 수 없습니다: {video_path}")

            yield cap

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"VideoCapture 생성 실패: {video_path}", e)
            raise
        finally:
            if cap is not None:
                cap.release()

    def get_video_info(self, video_path: Union[str, Path]) -> Dict[str, Any]:
        """비디오 정보 추출"""
        info = {
            'fps': 30.0,
            'total_frames': 0,
            'duration': 0.0,
            'width': 0,
            'height': 0,
            'codec': 'unknown'
        }

        try:
            with self.safe_video_capture(video_path) as cap:
                info['fps'] = cap.get(cv2.CAP_PROP_FPS) or 30.0
                info['total_frames'] = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
                info['width'] = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
                info['height'] = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)

                if info['fps'] > 0 and info['total_frames'] > 0:
                    info['duration'] = info['total_frames'] / info['fps']

                # 코덱 정보 (가능한 경우)
                fourcc = cap.get(cv2.CAP_PROP_FOURCC)
                if fourcc:
                    info['codec'] = "".join([chr((int(fourcc) >> 8 * i) & 0xFF) for i in range(4)])

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, f"비디오 정보 추출 실패: {video_path}", e)

        return info

    def calculate_sampling_params(self, video_info: Dict, config: ProcessingConfig) -> Dict:
        """샘플링 파라미터 계산"""
        fps = video_info['fps']
        duration = video_info['duration']
        total_frames = video_info['total_frames']

        # 기본 샘플링 간격 (프레임 단위)
        frame_interval = max(1, int(fps * config.frame_sampling_interval))

        # 최소 간격 적용
        min_frame_interval = max(1, int(fps * config.min_frame_interval))
        frame_interval = max(frame_interval, min_frame_interval)

        # 예상 프레임 수 계산
        estimated_frames = total_frames // frame_interval

        # 최대 프레임 수 제한
        if estimated_frames > config.max_frames_per_video:
            frame_interval = total_frames // config.max_frames_per_video
            frame_interval = max(frame_interval, min_frame_interval)

        return {
            'frame_interval': frame_interval,
            'estimated_frames': total_frames // frame_interval,
            'sampling_rate': frame_interval / fps
        }


# ============================================================================
# 유사도 계산 매니저
# ============================================================================

class SimilarityManager:
    """프레임 유사도 계산 관리"""

    def __init__(self):
        self.comparison_size = (640, 480)
        self.error_manager = ErrorManager()

    def calculate_ssim(self, img1: np.ndarray, img2: np.ndarray) -> float:
        """SSIM 기반 유사도 계산"""
        try:
            if img1.shape != img2.shape:
                img2 = cv2.resize(img2, (img1.shape[1], img1.shape[0]))

            gray1 = cv2.cvtColor(img1, cv2.COLOR_BGR2GRAY) if len(img1.shape) == 3 else img1
            gray2 = cv2.cvtColor(img2, cv2.COLOR_BGR2GRAY) if len(img2.shape) == 3 else img2

            if SSIM_AVAILABLE:
                similarity = ssim(gray1, gray2, data_range=gray1.max() - gray1.min())
            else:
                # MSE 기반 대체 계산
                diff = gray1.astype(np.float32) - gray2.astype(np.float32)
                mse = np.mean(diff ** 2)
                if mse == 0:
                    return 1.0
                similarity = max(0, 1.0 - (mse / (255.0 ** 2)))

            return float(similarity)

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "SSIM 계산 실패", e)
            return 0.0

    def calculate_histogram_similarity(self, img1: np.ndarray, img2: np.ndarray) -> float:
        """히스토그램 기반 유사도 계산"""
        try:
            hist1 = cv2.calcHist([img1], [0, 1, 2], None, [64, 64, 64], [0, 256, 0, 256, 0, 256])
            hist1 = cv2.normalize(hist1, hist1).flatten()

            hist2 = cv2.calcHist([img2], [0, 1, 2], None, [64, 64, 64], [0, 256, 0, 256, 0, 256])
            hist2 = cv2.normalize(hist2, hist2).flatten()

            similarity = cv2.compareHist(hist1, hist2, cv2.HISTCMP_CORREL)
            return float(max(0, similarity))

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "히스토그램 유사도 계산 실패", e)
            return 0.0

    def calculate_combined_similarity(self, img1: np.ndarray, img2: np.ndarray) -> float:
        """복합 유사도 계산"""
        weights = {'ssim': 0.6, 'histogram': 0.4}

        try:
            img1_resized = cv2.resize(img1, self.comparison_size)
            img2_resized = cv2.resize(img2, self.comparison_size)

            ssim_sim = self.calculate_ssim(img1_resized, img2_resized)
            hist_sim = self.calculate_histogram_similarity(img1_resized, img2_resized)

            combined = (ssim_sim * weights['ssim'] +
                        hist_sim * weights['histogram'])

            return float(combined)

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "복합 유사도 계산 실패", e)
            return 0.0


# ============================================================================
# 프레임 추출 매니저
# ============================================================================

class FrameExtractionManager:
    """프레임 추출 관리"""

    def __init__(self):
        self.video_manager = VideoManager()
        self.similarity_manager = SimilarityManager()
        self.path_manager = PathManager()
        self.error_manager = ErrorManager()
        self.config_manager = ConfigManager()

        self.frame_history = deque(maxlen=3)
        self.saved_frames = []

    def extract_frames(self, video_path: Union[str, Path], output_dir: Union[str, Path],
                       progress_callback=None) -> Tuple[List[str], int]:
        """프레임 추출 - 개선된 알고리즘"""

        video_path = self.path_manager.normalize_path(video_path)
        output_dir = self.path_manager.normalize_path(output_dir)

        if not video_path.exists():
            raise FileNotFoundError(f"비디오 파일이 없습니다: {video_path}")

        if not self.path_manager.ensure_directory(output_dir):
            raise OSError(f"출력 디렉토리 생성 실패: {output_dir}")

        config = self.config_manager.get_config()
        video_info = self.video_manager.get_video_info(video_path)
        sampling_params = self.video_manager.calculate_sampling_params(video_info, config)

        self.error_manager.log(LogLevel.INFO,
                               f"비디오 정보: FPS={video_info['fps']:.2f}, "
                               f"총 프레임={video_info['total_frames']}, "
                               f"샘플링 간격={sampling_params['frame_interval']}")

        self.saved_frames = []
        self.frame_history.clear()

        try:
            with self.video_manager.safe_video_capture(video_path) as cap:
                return self._extract_frames_internal(
                    cap, video_path, output_dir, video_info,
                    sampling_params, config, progress_callback
                )

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"프레임 추출 실패: {video_path}", e)
            return [], 0

    def _extract_frames_internal(self, cap, video_path: Path, output_dir: Path,
                                 video_info: Dict, sampling_params: Dict,
                                 config: ProcessingConfig, progress_callback) -> Tuple[List[str], int]:
        """내부 프레임 추출 로직"""

        frame_interval = sampling_params['frame_interval']
        total_frames = video_info['total_frames']
        fps = video_info['fps']

        frame_count = 0
        saved_count = 0
        recent_similarities = deque(maxlen=20)
        last_save_frame = -frame_interval

        video_name = self.path_manager.safe_filename(video_path.stem)

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            # 진행률 업데이트
            if progress_callback and total_frames > 0:
                progress = (frame_count / total_frames) * 100
                progress_callback(progress)

            # 샘플링 간격 확인
            if frame_count % frame_interval == 0:
                # 최소 간격 확인
                if frame_count - last_save_frame < frame_interval:
                    frame_count += 1
                    continue

                should_save, similarity = self._should_save_frame(frame, recent_similarities, config)

                if should_save:
                    success = self._save_frame(frame, frame_count, fps, output_dir, video_name)
                    if success:
                        saved_count += 1
                        last_save_frame = frame_count

                        # 히스토리 업데이트
                        frame_resized = cv2.resize(frame, self.similarity_manager.comparison_size)
                        self.frame_history.append(frame_resized.copy())

                        self.error_manager.log(LogLevel.INFO,
                                               f"프레임 저장 [{saved_count}]: {frame_count} "
                                               f"(유사도: {similarity:.3f})")

                if similarity > 0:
                    recent_similarities.append(similarity)

            frame_count += 1

            # 주기적 메모리 정리
            if frame_count % (frame_interval * 100) == 0:
                gc.collect()

        # 최소 프레임 보장
        if saved_count == 0 and total_frames > 0:
            self._save_fallback_frames(cap, total_frames, fps, output_dir, video_name)
            saved_count = len(self.saved_frames)

        self.error_manager.log(LogLevel.INFO,
                               f"프레임 추출 완료: {saved_count}개 저장 (총 {frame_count}개 처리)")

        return self.saved_frames, saved_count

    def _should_save_frame(self, current_frame: np.ndarray,
                           recent_similarities: deque,
                           config: ProcessingConfig) -> Tuple[bool, float]:
        """프레임 저장 여부 결정"""

        if not self.frame_history:
            return True, 0.0

        frame_resized = cv2.resize(current_frame, self.similarity_manager.comparison_size)

        similarities = []
        for hist_frame in self.frame_history:
            sim = self.similarity_manager.calculate_combined_similarity(hist_frame, frame_resized)
            similarities.append(sim)

        min_similarity = min(similarities)
        avg_similarity = np.mean(similarities)

        # 적응형 임계값 계산
        threshold = config.similarity_threshold
        if config.adaptive_threshold and recent_similarities:
            std_sim = np.std(recent_similarities)
            if std_sim < 0.05:  # 변화가 적음
                threshold -= 0.1
            elif std_sim > 0.15:  # 변화가 많음
                threshold += 0.05
            threshold = np.clip(threshold, 0.7, 0.98)

        should_save = min_similarity < threshold or avg_similarity < threshold * 0.95

        return should_save, avg_similarity

    def _save_frame(self, frame: np.ndarray, frame_count: int, fps: float,
                    output_dir: Path, video_name: str) -> bool:
        """프레임 저장"""
        try:
            second = int(frame_count / max(1, fps))
            frame_filename = output_dir / f"{video_name}_frame_{second:04d}s_{frame_count:06d}.jpg"

            if self.path_manager.write_image_korean(frame, frame_filename, quality=90):
                self.saved_frames.append(str(frame_filename))
                return True

            return False

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"프레임 저장 실패: {frame_count}", e)
            return False

    def _save_fallback_frames(self, cap, total_frames: int, fps: float,
                              output_dir: Path, video_name: str):
        """최소 프레임 보장을 위한 폴백 저장"""
        try:
            self.error_manager.log(LogLevel.WARNING, "최소 프레임 보장을 위해 강제 저장합니다.")

            sample_positions = np.linspace(0, total_frames - 1, min(5, total_frames), dtype=int)

            for pos in sample_positions:
                cap.set(cv2.CAP_PROP_POS_FRAMES, pos)
                ret, frame = cap.read()
                if ret:
                    if self._save_frame(frame, pos, fps, output_dir, video_name):
                        pass  # 이미 saved_frames에 추가됨

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "폴백 프레임 저장 실패", e)


# ============================================================================
# 개선된 GUI 클래스
# ============================================================================

class VideoProcessorGUI:
    """개선된 GUI 클래스"""

    def __init__(self, root):
        self.root = root
        self.root.title("🎬 올인원 비디오 처리기 v4.0 - 리팩토링 버전")
        self.root.geometry("1000x900")

        self.config_manager = ConfigManager()
        self.error_manager = ErrorManager()
        self.path_manager = PathManager()

        # OCRManager 초기화
        if OCR_AVAILABLE:
            try:
                from ocr_manager import OCRManager
                self.ocr_manager = OCRManager(use_ai_model=False)
                self.error_manager.log(LogLevel.INFO, "OCRManager 초기화 완료")
            except Exception as e:
                self.ocr_manager = None
                self.error_manager.log(LogLevel.WARNING, f"OCRManager 초기화 실패: {e}")
        else:
            self.ocr_manager = None

        self.video_files = []
        self.processing = False
        self.processing_thread = None
        self.log_queue = queue.Queue(maxsize=1000)

        self.setup_gui()
        self.setup_logging()
        self.load_gui_config()
        self.update_log()

        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)

    def setup_gui(self):
        """GUI 레이아웃 설정"""
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))

        # 1. 파일 선택 섹션
        self.create_file_selection_section(main_frame, 0)

        # 2. 처리 옵션 섹션
        self.create_processing_options_section(main_frame, 1)

        # 3. 샘플링 설정 섹션 (새로 추가)
        self.create_sampling_section(main_frame, 2)

        # 4. 출력 설정 섹션
        self.create_output_section(main_frame, 3)

        # 5. 진행 상황 섹션
        self.create_progress_section(main_frame, 4)

        # 6. 로그 섹션
        self.create_log_section(main_frame, 5)

        # 7. 컨트롤 버튼 섹션
        self.create_control_section(main_frame, 6)

        # 그리드 설정
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(0, weight=1)
        main_frame.rowconfigure(5, weight=1)  # 로그 섹션만 확장

    def create_file_selection_section(self, parent, row):
        """파일 선택 섹션"""
        file_frame = ttk.LabelFrame(parent, text="📁 비디오 파일 선택", padding="10")
        file_frame.grid(row=row, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)

        # 버튼들
        button_frame = ttk.Frame(file_frame)
        button_frame.grid(row=0, column=0, sticky=(tk.W, tk.E))

        ttk.Button(button_frame, text="파일 선택", command=self.select_files).grid(row=0, column=0, padx=5)
        ttk.Button(button_frame, text="폴더 선택", command=self.select_folder).grid(row=0, column=1, padx=5)
        ttk.Button(button_frame, text="선택 제거", command=self.remove_selected).grid(row=0, column=2, padx=5)
        ttk.Button(button_frame, text="모두 제거", command=self.clear_files).grid(row=0, column=3, padx=5)

        # 파일 목록
        list_frame = ttk.Frame(file_frame)
        list_frame.grid(row=1, column=0, pady=5, sticky=(tk.W, tk.E))

        self.file_listbox = tk.Listbox(list_frame, height=4, width=70)
        scrollbar = ttk.Scrollbar(list_frame, orient="vertical", command=self.file_listbox.yview)
        self.file_listbox.configure(yscrollcommand=scrollbar.set)

        self.file_listbox.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        scrollbar.grid(row=0, column=1, sticky=(tk.N, tk.S))

        list_frame.columnconfigure(0, weight=1)
        file_frame.columnconfigure(0, weight=1)

    def create_processing_options_section(self, parent, row):
        """처리 옵션 섹션"""
        options_frame = ttk.LabelFrame(parent, text="⚙️ 처리 옵션", padding="10")
        options_frame.grid(row=row, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)

        # 좌측: 프레임 추출 옵션
        frame_options = ttk.LabelFrame(options_frame, text="📷 프레임 추출", padding="5")
        frame_options.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N), padx=5, pady=5)

        self.extract_frames_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(frame_options, text="프레임 추출",
                        variable=self.extract_frames_var).grid(row=0, column=0, columnspan=3, sticky=tk.W)

        ttk.Label(frame_options, text="유사도 임계값:").grid(row=1, column=0, sticky=tk.W, padx=(20, 0))
        self.threshold_var = tk.DoubleVar(value=0.95)
        self.threshold_scale = ttk.Scale(frame_options, from_=0.5, to=1.0,
                                         variable=self.threshold_var, orient="horizontal", length=150)
        self.threshold_scale.grid(row=1, column=1, padx=5)
        self.threshold_label = ttk.Label(frame_options, text="0.95")
        self.threshold_label.grid(row=1, column=2)
        self.threshold_scale.configure(command=self._update_threshold_label)

        self.adaptive_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(frame_options, text="적응형 임계값",
                        variable=self.adaptive_var).grid(row=2, column=0, columnspan=3, sticky=tk.W, padx=(20, 0))

        # 우측: 출력 옵션
        output_options = ttk.LabelFrame(options_frame, text="📊 출력 옵션", padding="5")
        output_options.grid(row=0, column=1, sticky=(tk.W, tk.E, tk.N), padx=5, pady=5)

        self.create_ppt_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(output_options, text="PPT 생성",
                        variable=self.create_ppt_var).grid(row=0, column=0, columnspan=3, sticky=tk.W)

        ttk.Label(output_options, text="그리드:").grid(row=1, column=0, sticky=tk.W, padx=(20, 0))
        grid_frame = ttk.Frame(output_options)
        grid_frame.grid(row=1, column=1, columnspan=2, padx=5)

        self.grid_rows_var = tk.IntVar(value=3)
        self.grid_cols_var = tk.IntVar(value=3)

        ttk.Spinbox(grid_frame, from_=1, to=10, width=5,
                    textvariable=self.grid_rows_var).grid(row=0, column=0)
        ttk.Label(grid_frame, text=" × ").grid(row=0, column=1)
        ttk.Spinbox(grid_frame, from_=1, to=10, width=5,
                    textvariable=self.grid_cols_var).grid(row=0, column=2)

        self.extract_audio_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(output_options, text="음성 추출",
                        variable=self.extract_audio_var).grid(row=2, column=0, columnspan=3, sticky=tk.W)

        self.create_srt_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(output_options, text="SRT 자막",
                        variable=self.create_srt_var).grid(row=3, column=0, columnspan=3, sticky=tk.W, padx=(20, 0))

        self.create_word_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(output_options, text="Word 보고서",
                        variable=self.create_word_var).grid(row=4, column=0, columnspan=3, sticky=tk.W, padx=(20, 0))

        # OCR 오탈자 검수 (선택적 - 시간 소요)
        self.ocr_spell_check_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(output_options, text="OCR 오탈자 검수 (느림)",
                        variable=self.ocr_spell_check_var).grid(row=5, column=0, columnspan=3, sticky=tk.W, padx=(20, 0))

        # Whisper 모델 선택
        ttk.Label(output_options, text="Whisper 모델:").grid(row=6, column=0, sticky=tk.W, padx=(20, 0))
        self.whisper_model_var = tk.StringVar(value="base")
        ttk.Combobox(output_options, textvariable=self.whisper_model_var,
                     values=["tiny", "base", "small", "medium", "large"],
                     state="readonly", width=10).grid(row=6, column=1, padx=5)

        options_frame.columnconfigure(0, weight=1)
        options_frame.columnconfigure(1, weight=1)

    def create_sampling_section(self, parent, row):
        """샘플링 설정 섹션 - 새로 추가"""
        sampling_frame = ttk.LabelFrame(parent, text="🎯 샘플링 설정", padding="10")
        sampling_frame.grid(row=row, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)

        # 샘플링 간격 설정
        ttk.Label(sampling_frame, text="샘플링 간격 (초):").grid(row=0, column=0, sticky=tk.W)
        self.sampling_interval_var = tk.DoubleVar(value=1.0)
        self.sampling_scale = ttk.Scale(sampling_frame, from_=0.1, to=10.0,
                                        variable=self.sampling_interval_var,
                                        orient="horizontal", length=200)
        self.sampling_scale.grid(row=0, column=1, padx=5)
        self.sampling_label = ttk.Label(sampling_frame, text="1.0초")
        self.sampling_label.grid(row=0, column=2)
        self.sampling_scale.configure(command=self._update_sampling_label)

        # 최소 간격 설정
        ttk.Label(sampling_frame, text="최소 간격 (초):").grid(row=1, column=0, sticky=tk.W)
        self.min_interval_var = tk.DoubleVar(value=0.5)
        self.min_interval_scale = ttk.Scale(sampling_frame, from_=0.1, to=5.0,
                                            variable=self.min_interval_var,
                                            orient="horizontal", length=200)
        self.min_interval_scale.grid(row=1, column=1, padx=5)
        self.min_interval_label = ttk.Label(sampling_frame, text="0.5초")
        self.min_interval_label.grid(row=1, column=2)
        self.min_interval_scale.configure(command=self._update_min_interval_label)

        # 최대 프레임 수 설정
        ttk.Label(sampling_frame, text="최대 프레임 수:").grid(row=2, column=0, sticky=tk.W)
        self.max_frames_var = tk.IntVar(value=200)
        ttk.Spinbox(sampling_frame, from_=10, to=1000, width=10,
                    textvariable=self.max_frames_var).grid(row=2, column=1, padx=5, sticky=tk.W)

        # 미리보기 정보
        self.preview_label = ttk.Label(sampling_frame, text="예상 프레임 수: 계산 중...",
                                       foreground="blue")
        self.preview_label.grid(row=3, column=0, columnspan=3, pady=5)

        # 설정 변경 시 미리보기 업데이트
        for var in [self.sampling_interval_var, self.min_interval_var, self.max_frames_var]:
            var.trace('w', self._update_preview)

    def create_output_section(self, parent, row):
        """출력 설정 섹션"""
        output_frame = ttk.LabelFrame(parent, text="📂 출력 설정", padding="10")
        output_frame.grid(row=row, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)

        ttk.Label(output_frame, text="출력 폴더:").grid(row=0, column=0)
        self.output_dir_var = tk.StringVar(value="output")
        ttk.Entry(output_frame, textvariable=self.output_dir_var, width=60).grid(row=0, column=1, padx=5)
        ttk.Button(output_frame, text="찾아보기", command=self.select_output_dir).grid(row=0, column=2)

        output_frame.columnconfigure(1, weight=1)

    def create_progress_section(self, parent, row):
        """진행 상황 섹션"""
        progress_frame = ttk.LabelFrame(parent, text="📈 진행 상황", padding="10")
        progress_frame.grid(row=row, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=5)

        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(progress_frame, variable=self.progress_var,
                                            maximum=100, length=500)
        self.progress_bar.grid(row=0, column=0, columnspan=2, pady=5)

        self.status_label = ttk.Label(progress_frame, text="대기 중...")
        self.status_label.grid(row=1, column=0, columnspan=2)

        # 상태 정보 표시
        info_frame = ttk.Frame(progress_frame)
        info_frame.grid(row=2, column=0, columnspan=2, pady=5)

        self.current_file_label = ttk.Label(info_frame, text="")
        self.current_file_label.grid(row=0, column=0)

        self.eta_label = ttk.Label(info_frame, text="")
        self.eta_label.grid(row=1, column=0)

    def create_log_section(self, parent, row):
        """로그 섹션"""
        log_frame = ttk.LabelFrame(parent, text="📝 로그", padding="10")
        log_frame.grid(row=row, column=0, columnspan=2, sticky=(tk.W, tk.E, tk.N, tk.S), pady=5)

        self.log_text = scrolledtext.ScrolledText(log_frame, width=90, height=12)
        self.log_text.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))

        log_frame.columnconfigure(0, weight=1)
        log_frame.rowconfigure(0, weight=1)

    def create_control_section(self, parent, row):
        """컨트롤 버튼 섹션"""
        button_frame = ttk.Frame(parent)
        button_frame.grid(row=row, column=0, columnspan=2, pady=10)

        self.process_button = ttk.Button(button_frame, text="🚀 처리 시작",
                                         command=self.start_processing)
        self.process_button.grid(row=0, column=0, padx=5)

        self.stop_button = ttk.Button(button_frame, text="⏹️ 중지",
                                      command=self.stop_processing, state="disabled")
        self.stop_button.grid(row=0, column=1, padx=5)

        ttk.Button(button_frame, text="💾 설정 저장",
                   command=self.save_config).grid(row=0, column=2, padx=5)

        ttk.Button(button_frame, text="🔄 설정 불러오기",
                   command=self.load_config).grid(row=0, column=3, padx=5)

        # OCR 검수만 실행 버튼
        ttk.Button(button_frame, text="🔍 OCR 검수만",
                   command=self.run_ocr_only).grid(row=0, column=4, padx=5)

        # 통합 맞춤법 검사 버튼 (새로 추가)
        ttk.Button(button_frame, text="📝 맞춤법 검사",
                   command=self.run_integrated_spell_check).grid(row=0, column=5, padx=5)

        ttk.Button(button_frame, text="🗑️ 로그 지우기",
                   command=self.clear_log).grid(row=0, column=6, padx=5)

        ttk.Button(button_frame, text="📊 상태 확인",
                   command=self.show_status).grid(row=0, column=7, padx=5)

        ttk.Button(button_frame, text="❌ 종료",
                   command=self.on_closing).grid(row=1, column=0, columnspan=8, padx=5, pady=(5,0))

    def _update_threshold_label(self, value):
        """유사도 임계값 라벨 업데이트"""
        self.threshold_label.configure(text=f"{float(value):.2f}")

    def _update_sampling_label(self, value):
        """샘플링 간격 라벨 업데이트"""
        self.sampling_label.configure(text=f"{float(value):.1f}초")
        self._update_preview()

    def _update_min_interval_label(self, value):
        """최소 간격 라벨 업데이트"""
        self.min_interval_label.configure(text=f"{float(value):.1f}초")
        self._update_preview()

    def _update_preview(self, *args):
        """예상 프레임 수 미리보기 업데이트"""
        try:
            if not self.video_files:
                self.preview_label.configure(text="비디오 파일을 선택하세요")
                return

            # 첫 번째 비디오로 예상치 계산
            first_video = Path(self.video_files[0])
            if first_video.exists():
                video_manager = VideoManager()
                video_info = video_manager.get_video_info(first_video)

                fps = video_info.get('fps', 30)
                duration = video_info.get('duration', 0)

                if duration > 0:
                    sampling_interval = self.sampling_interval_var.get()
                    estimated_frames = int(duration / sampling_interval)
                    max_frames = self.max_frames_var.get()

                    actual_frames = min(estimated_frames, max_frames)

                    self.preview_label.configure(
                        text=f"예상 프레임 수: {actual_frames}개 (영상 길이: {duration:.1f}초)"
                    )
                else:
                    self.preview_label.configure(text="비디오 정보를 읽을 수 없습니다")
            else:
                self.preview_label.configure(text="비디오 파일이 존재하지 않습니다")

        except Exception:
            self.preview_label.configure(text="미리보기 계산 실패")

    def select_files(self):
        """파일 선택"""
        try:
            files = filedialog.askopenfilenames(
                title="비디오 파일 선택",
                filetypes=[
                    ("비디오 파일", "*.mp4 *.avi *.mov *.mkv *.wmv *.flv *.webm *.m4v"),
                    ("모든 파일", "*.*")
                ]
            )

            added_count = 0
            for file in files:
                if file not in self.video_files:
                    self.video_files.append(file)
                    # 한글 파일명도 제대로 표시
                    display_name = Path(file).name
                    self.file_listbox.insert(tk.END, display_name)
                    added_count += 1

            if added_count > 0:
                self.error_manager.log(LogLevel.INFO, f"{added_count}개 파일이 추가되었습니다.")
                self._update_preview()

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "파일 선택 실패", e)
            messagebox.showerror("오류", f"파일 선택 중 오류가 발생했습니다:\n{str(e)}")

    def select_folder(self):
        """폴더 선택"""
        try:
            folder = filedialog.askdirectory(title="폴더 선택")
            if folder:
                folder_path = Path(folder)
                video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v'}
                added_count = 0

                for file_path in folder_path.rglob('*'):
                    if file_path.suffix.lower() in video_extensions and file_path.is_file():
                        file_str = str(file_path)
                        if file_str not in self.video_files:
                            self.video_files.append(file_str)
                            self.file_listbox.insert(tk.END, file_path.name)
                            added_count += 1

                if added_count > 0:
                    self.error_manager.log(LogLevel.INFO, f"폴더에서 {added_count}개 비디오 파일이 추가되었습니다.")
                    self._update_preview()
                else:
                    messagebox.showinfo("정보", "폴더에서 비디오 파일을 찾을 수 없습니다.")

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "폴더 선택 실패", e)
            messagebox.showerror("오류", f"폴더 선택 중 오류가 발생했습니다:\n{str(e)}")

    def remove_selected(self):
        """선택된 파일 제거"""
        try:
            selections = self.file_listbox.curselection()
            if selections:
                for index in reversed(selections):
                    self.file_listbox.delete(index)
                    if 0 <= index < len(self.video_files):
                        del self.video_files[index]
                self.error_manager.log(LogLevel.INFO, f"{len(selections)}개 파일이 제거되었습니다.")
                self._update_preview()
        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "파일 제거 실패", e)

    def clear_files(self):
        """모든 파일 제거"""
        try:
            self.file_listbox.delete(0, tk.END)
            count = len(self.video_files)
            self.video_files.clear()
            if count > 0:
                self.error_manager.log(LogLevel.INFO, f"{count}개 파일이 모두 제거되었습니다.")
                self._update_preview()
        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "파일 목록 지우기 실패", e)

    def select_output_dir(self):
        """출력 디렉토리 선택"""
        try:
            directory = filedialog.askdirectory(title="출력 폴더 선택")
            if directory:
                self.output_dir_var.set(directory)
        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "출력 디렉토리 선택 실패", e)

    def save_config(self):
        """현재 GUI 설정을 저장"""
        try:
            self.config_manager.update_config(
                similarity_threshold=self.threshold_var.get(),
                adaptive_threshold=self.adaptive_var.get(),
                extract_frames=self.extract_frames_var.get(),
                create_ppt=self.create_ppt_var.get(),
                extract_audio=self.extract_audio_var.get(),
                create_srt=self.create_srt_var.get(),
                create_word=self.create_word_var.get(),
                ocr_spell_check=self.ocr_spell_check_var.get(),
                grid_rows=self.grid_rows_var.get(),
                grid_cols=self.grid_cols_var.get(),
                whisper_model=self.whisper_model_var.get(),
                frame_sampling_interval=self.sampling_interval_var.get(),
                min_frame_interval=self.min_interval_var.get(),
                max_frames_per_video=self.max_frames_var.get()
            )
            messagebox.showinfo("설정 저장", "설정이 저장되었습니다.")

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "설정 저장 실패", e)
            messagebox.showerror("오류", f"설정 저장 실패:\n{str(e)}")

    def load_config(self):
        """저장된 설정을 불러오기"""
        try:
            self.config_manager.load_config()
            self.load_gui_config()
            messagebox.showinfo("설정 불러오기", "설정이 불러와졌습니다.")

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "설정 불러오기 실패", e)
            messagebox.showerror("오류", f"설정 불러오기 실패:\n{str(e)}")

    def load_gui_config(self):
        """설정을 GUI에 반영"""
        try:
            config = self.config_manager.get_config()

            self.threshold_var.set(config.similarity_threshold)
            self.adaptive_var.set(config.adaptive_threshold)
            self.extract_frames_var.set(config.extract_frames)
            self.create_ppt_var.set(config.create_ppt)
            self.extract_audio_var.set(config.extract_audio)
            self.create_srt_var.set(config.create_srt)
            self.create_word_var.set(config.create_word)
            self.ocr_spell_check_var.set(config.ocr_spell_check)
            self.grid_rows_var.set(config.grid_rows)
            self.grid_cols_var.set(config.grid_cols)
            self.whisper_model_var.set(config.whisper_model)
            self.sampling_interval_var.set(config.frame_sampling_interval)
            self.min_interval_var.set(config.min_frame_interval)
            self.max_frames_var.set(config.max_frames_per_video)

            # 라벨 업데이트
            self._update_threshold_label(config.similarity_threshold)
            self._update_sampling_label(config.frame_sampling_interval)
            self._update_min_interval_label(config.min_frame_interval)

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "GUI 설정 로드 실패", e)

    def clear_log(self):
        """로그 지우기"""
        try:
            self.log_text.delete(1.0, tk.END)
            while not self.log_queue.empty():
                try:
                    self.log_queue.get_nowait()
                except queue.Empty:
                    break
        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "로그 지우기 실패", e)

    def show_status(self):
        """현재 상태 표시"""
        try:
            config = self.config_manager.get_config()

            status_info = [
                f"선택된 파일: {len(self.video_files)}개",
                f"처리 중: {'예' if self.processing else '아니오'}",
                f"출력 폴더: {self.output_dir_var.get()}",
                "",
                "=== 처리 설정 ===",
                f"유사도 임계값: {config.similarity_threshold:.2f}",
                f"적응형 임계값: {'사용' if config.adaptive_threshold else '미사용'}",
                f"샘플링 간격: {config.frame_sampling_interval}초",
                f"최소 간격: {config.min_frame_interval}초",
                f"최대 프레임 수: {config.max_frames_per_video}개",
                f"그리드 크기: {config.grid_rows}x{config.grid_cols}",
                f"Whisper 모델: {config.whisper_model}",
                "",
                "=== 라이브러리 상태 ===",
                f"OpenCV: {'사용 가능' if CV2_AVAILABLE else '없음'}",
                f"SSIM: {'사용 가능' if SSIM_AVAILABLE else '기본 방법 사용'}",
                f"PPT: {'사용 가능' if PPTX_AVAILABLE else '없음'}",
                f"Whisper: {'사용 가능' if WHISPER_AVAILABLE else '없음'}",
                f"MoviePy: {'사용 가능' if MOVIEPY_AVAILABLE else '없음'}",
                f"Word: {'사용 가능' if DOCX_AVAILABLE else '없음'}"
            ]

            messagebox.showinfo("현재 상태", "\n".join(status_info))

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "상태 표시 실패", e)

    def validate_inputs(self) -> List[str]:
        """입력값 검증"""
        errors = []

        try:
            if not self.video_files:
                errors.append("처리할 비디오 파일을 선택하세요")

            threshold = self.threshold_var.get()
            if not (0.0 <= threshold <= 1.0):
                errors.append(f"유사도 임계값은 0.0~1.0 사이여야 합니다: {threshold}")

            sampling_interval = self.sampling_interval_var.get()
            min_interval = self.min_interval_var.get()
            if sampling_interval < min_interval:
                errors.append("샘플링 간격은 최소 간격보다 커야 합니다")

            rows, cols = self.grid_rows_var.get(), self.grid_cols_var.get()
            if not (1 <= rows <= 10 and 1 <= cols <= 10):
                errors.append(f"그리드 크기는 1~10 사이여야 합니다: {rows}x{cols}")

            output_dir = self.output_dir_var.get()
            if not output_dir:
                errors.append("출력 폴더를 설정하세요")
            else:
                try:
                    if not self.path_manager.ensure_directory(output_dir):
                        errors.append("출력 폴더 생성 실패")
                except Exception as e:
                    errors.append(f"출력 폴더 접근 실패: {e}")

            # 파일 존재 확인
            missing_files = []
            for video_path in self.video_files:
                if not Path(video_path).exists():
                    missing_files.append(Path(video_path).name)

            if missing_files:
                if len(missing_files) <= 3:
                    errors.append(f"존재하지 않는 파일: {', '.join(missing_files)}")
                else:
                    errors.append(f"존재하지 않는 파일 {len(missing_files)}개: {', '.join(missing_files[:3])}...")

        except Exception as e:
            errors.append(f"입력값 검증 중 오류: {e}")

        return errors

    def setup_logging(self):
        """로깅 설정"""

        class QueueHandler(logging.Handler):
            def __init__(self, queue):
                super().__init__()
                self.queue = queue

            def emit(self, record):
                try:
                    if not self.queue.full():
                        self.queue.put(self.format(record))
                except Exception:
                    pass

        queue_handler = QueueHandler(self.log_queue)
        queue_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))

        logger = logging.getLogger('VideoProcessor')
        logger.addHandler(queue_handler)
        logger.setLevel(logging.INFO)

    def update_log(self):
        """로그 업데이트"""
        try:
            message_count = 0
            while not self.log_queue.empty() and message_count < 10:
                try:
                    message = self.log_queue.get_nowait()
                    self.log_text.insert(tk.END, message + '\n')
                    message_count += 1
                except queue.Empty:
                    break

            if message_count > 0:
                self.log_text.see(tk.END)

        except Exception:
            pass

        self.root.after(100, self.update_log)

    def start_processing(self):
        """비디오 처리 시작"""
        if self.processing:
            messagebox.showinfo("정보", "이미 처리 중입니다.")
            return

        errors = self.validate_inputs()
        if errors:
            messagebox.showwarning("입력 오류", "다음 문제를 해결하세요:\n\n" + "\n".join(errors))
            return

        try:
            # GUI 설정을 ConfigManager에 반영
            self.config_manager.update_config(
                similarity_threshold=self.threshold_var.get(),
                adaptive_threshold=self.adaptive_var.get(),
                extract_frames=self.extract_frames_var.get(),
                create_ppt=self.create_ppt_var.get(),
                extract_audio=self.extract_audio_var.get(),
                create_srt=self.create_srt_var.get(),
                create_word=self.create_word_var.get(),
                ocr_spell_check=self.ocr_spell_check_var.get(),
                grid_rows=self.grid_rows_var.get(),
                grid_cols=self.grid_cols_var.get(),
                whisper_model=self.whisper_model_var.get(),
                frame_sampling_interval=self.sampling_interval_var.get(),
                min_frame_interval=self.min_interval_var.get(),
                max_frames_per_video=self.max_frames_var.get()
            )

            self.processing = True
            self.process_button.configure(state="disabled")
            self.stop_button.configure(state="normal")
            self.progress_var.set(0)

            self.processing_thread = threading.Thread(target=self.process_videos, daemon=True)
            self.processing_thread.start()

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "처리 시작 실패", e)
            self.processing = False
            self.process_button.configure(state="normal")
            self.stop_button.configure(state="disabled")

    def stop_processing(self):
        """처리 중지"""
        try:
            self.processing = False
            self._safe_gui_update(lambda: self.status_label.configure(text="중지 중..."))
            self.error_manager.log(LogLevel.INFO, "사용자가 처리 중지를 요청했습니다.")
        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "처리 중지 실패", e)

    def run_ocr_only(self):
        """저장된 이미지에 대해 OCR 검수만 실행"""
        try:
            from tkinter import filedialog, messagebox

            # 입력 폴더 선택
            input_folder = filedialog.askdirectory(title="OCR 검수할 이미지 폴더 선택")
            if not input_folder:
                return

            input_path = Path(input_folder)
            if not input_path.exists():
                messagebox.showerror("오류", f"폴더가 존재하지 않습니다:\n{input_folder}")
                return

            # 출력 폴더 자동 설정
            output_path = input_path.parent / f"{input_path.name}_ocr_checked"

            # 확인 다이얼로그
            msg = f"OCR 검수를 시작하시겠습니까?\n\n"
            msg += f"입력 폴더: {input_path}\n"
            msg += f"출력 폴더: {output_path}\n\n"
            msg += f"비교 모드: 빨간색(오류) vs 초록색(교정)"

            if not messagebox.askyesno("OCR 검수 확인", msg):
                return

            # OCR Manager 확인
            if not self.ocr_manager or not OCR_AVAILABLE:
                messagebox.showerror("오류", "OCR Manager를 사용할 수 없습니다.\nocr_manager.py를 확인하세요.")
                return

            # 로그 시작
            self.error_manager.log(LogLevel.INFO, "=" * 60)
            self.error_manager.log(LogLevel.INFO, "OCR 검수만 실행")
            self.error_manager.log(LogLevel.INFO, f"입력: {input_path}")
            self.error_manager.log(LogLevel.INFO, f"출력: {output_path}")
            self.error_manager.log(LogLevel.INFO, "=" * 60)

            # 진행 콜백
            def progress_callback(current, total, path):
                progress = (current / total) * 100
                self._safe_gui_update(lambda: self.progress_var.set(progress))
                self._safe_gui_update(lambda: self.status_label.configure(
                    text=f"OCR 검수 중: {current}/{total}"))
                self.error_manager.log(LogLevel.INFO,
                    f"[{current}/{total}] {Path(path).name}")

            # OCR 처리 시작
            self._safe_gui_update(lambda: self.process_button.configure(state="disabled"))

            result = self.ocr_manager.process_folder(
                input_path,
                output_dir=output_path,
                comparison_mode=True,
                file_pattern="*.jpg",  # 필요시 "*.png"로 변경
                callback=progress_callback
            )

            # 결과 표시
            self.error_manager.log(LogLevel.INFO, "=" * 60)
            if result['success']:
                self.error_manager.log(LogLevel.INFO,
                    f"✅ OCR 검수 완료: {result['processed']}/{result['total']}개")
                self.error_manager.log(LogLevel.INFO, f"출력 폴더: {output_path}")

                messagebox.showinfo("완료",
                    f"OCR 검수가 완료되었습니다!\n\n"
                    f"처리: {result['processed']}/{result['total']}개\n"
                    f"실패: {result['failed']}개\n\n"
                    f"출력: {output_path}")
            else:
                self.error_manager.log(LogLevel.ERROR, "❌ OCR 검수 실패")
                messagebox.showerror("실패", "OCR 검수 중 오류가 발생했습니다.")

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "OCR 검수 중 오류 발생", e)
            import traceback
            traceback.print_exc()
            messagebox.showerror("오류", f"OCR 검수 중 오류:\n{e}")

        finally:
            self._safe_gui_update(lambda: self.process_button.configure(state="normal"))
            self._safe_gui_update(lambda: self.progress_var.set(0))
            self._safe_gui_update(lambda: self.status_label.configure(text="대기 중"))

    def run_integrated_spell_check(self):
        """통합 맞춤법 검사 실행 (문서 → 사전 생성 → 맞춤법 검사 → 리포트)"""
        try:
            from tkinter import filedialog, messagebox

            # 입력 폴더 선택
            input_folder = filedialog.askdirectory(title="맞춤법 검사할 문서 폴더 선택 (Word/PPT)")
            if not input_folder:
                return

            input_path = Path(input_folder)
            if not input_path.exists():
                messagebox.showerror("오류", f"폴더가 존재하지 않습니다:\n{input_folder}")
                return

            # 출력 폴더 자동 설정
            output_path = input_path.parent / f"{input_path.name}_spell_checked"

            # 확인 다이얼로그
            msg = f"통합 맞춤법 검사를 시작하시겠습니까?\n\n"
            msg += f"입력 폴더: {input_path}\n"
            msg += f"출력 폴더: {output_path}\n\n"
            msg += f"기능:\n"
            msg += f"  1. Word/PPT 문서 텍스트 추출\n"
            msg += f"  2. 자주 사용되는 단어 사전 생성\n"
            msg += f"  3. 사전 단어를 화이트리스트로 반영\n"
            msg += f"  4. 실제 오탈자만 검출\n"
            msg += f"  5. 엑셀 리포트 생성"

            if not messagebox.askyesno("통합 맞춤법 검사 확인", msg):
                return

            # 통합 솔루션 import
            try:
                from integrated_spell_solution import IntegratedSpellSolution
            except ImportError:
                messagebox.showerror("오류", "integrated_spell_solution.py를 찾을 수 없습니다.")
                return

            # 로그 시작
            self.error_manager.log(LogLevel.INFO, "=" * 60)
            self.error_manager.log(LogLevel.INFO, "통합 맞춤법 검사 실행")
            self.error_manager.log(LogLevel.INFO, f"입력: {input_path}")
            self.error_manager.log(LogLevel.INFO, f"출력: {output_path}")
            self.error_manager.log(LogLevel.INFO, "=" * 60)

            # 진행 콜백
            def progress_callback(stage_name, current, total):
                progress = (current / total) * 100
                self._safe_gui_update(lambda: self.progress_var.set(progress))
                self._safe_gui_update(lambda: self.status_label.configure(
                    text=f"[{stage_name}] 진행 중..."))
                self.error_manager.log(LogLevel.INFO, f"[{stage_name}] 진행 중...")

            # 처리 시작
            self._safe_gui_update(lambda: self.process_button.configure(state="disabled"))

            # 통합 솔루션 실행
            solution = IntegratedSpellSolution(
                use_morpheme=True,       # 형태소 분석 사용
                morpheme_engine='okt',   # Okt 사용 (빠름)
                use_ai_spell_check=False  # AI 모델 비활성화 (빠름)
            )

            result = solution.run_integrated_pipeline(
                input_folder=str(input_path),
                output_dir=str(output_path),
                file_patterns=['*.doc', '*.docx', '*.ppt', '*.pptx'],
                password=None,
                top_k_words=500,
                min_frequency=2,
                min_priority=0.05,
                callback=progress_callback
            )

            # 결과 표시
            self.error_manager.log(LogLevel.INFO, "=" * 60)
            if result['success']:
                self.error_manager.log(LogLevel.INFO, "✅ 통합 맞춤법 검사 완료")
                self.error_manager.log(LogLevel.INFO, f"추출 파일: {result['extracted_files']}개")
                self.error_manager.log(LogLevel.INFO, f"전체 단어: {result['total_words']}개")
                self.error_manager.log(LogLevel.INFO, f"사전 단어: {result['dictionary_words']}개")
                self.error_manager.log(LogLevel.INFO, f"오타 발견: {result['typos_found']}개")
                self.error_manager.log(LogLevel.INFO, f"출력 폴더: {output_path}")

                # 종합 리포트 경로
                report_path = result['output_files'].get('comprehensive_report', '')

                messagebox.showinfo("완료",
                    f"통합 맞춤법 검사가 완료되었습니다!\n\n"
                    f"추출 파일: {result['extracted_files']}개\n"
                    f"전체 단어: {result['total_words']}개\n"
                    f"사전 단어: {result['dictionary_words']}개 (화이트리스트)\n"
                    f"오타 발견: {result['typos_found']}개\n"
                    f"이미지 추출: {result.get('images_extracted', 0)}개\n"
                    f"OCR 처리: {result.get('ocr_images_processed', 0)}개\n\n"
                    f"출력 폴더: {output_path}\n\n"
                    f"📊 종합 리포트:\n{report_path}")
            else:
                self.error_manager.log(LogLevel.ERROR, "❌ 통합 맞춤법 검사 실패")
                error_msg = result.get('error', '알 수 없는 오류')
                messagebox.showerror("실패", f"맞춤법 검사 중 오류 발생:\n{error_msg}")

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "통합 맞춤법 검사 중 오류 발생", e)
            import traceback
            traceback.print_exc()
            messagebox.showerror("오류", f"맞춤법 검사 중 오류:\n{e}")

        finally:
            self._safe_gui_update(lambda: self.process_button.configure(state="normal"))
            self._safe_gui_update(lambda: self.progress_var.set(0))
            self._safe_gui_update(lambda: self.status_label.configure(text="대기 중"))

    def _safe_gui_update(self, func):
        """스레드 안전한 GUI 업데이트"""
        self.root.after(0, func)

    def process_videos(self):
        """비디오 처리 (별도 스레드)"""
        try:
            processor = VideoProcessorMain()

            total_videos = len(self.video_files)
            completed_videos = 0
            start_time = time.time()

            for i, video_path in enumerate(self.video_files):
                if not self.processing:
                    break

                video_name = Path(video_path).name
                self._safe_gui_update(lambda name=video_name:
                                      self.status_label.configure(text=f"처리 중: {name}"))
                self._safe_gui_update(lambda name=video_name:
                                      self.current_file_label.configure(text=f"현재 파일: {name}"))

                def progress_callback(progress):
                    overall_progress = ((i / total_videos) + (progress / 100 / total_videos)) * 100
                    self._safe_gui_update(lambda p=overall_progress: self.progress_var.set(p))

                    # ETA 계산
                    if progress > 0:
                        elapsed = time.time() - start_time
                        total_estimated = elapsed / (overall_progress / 100)
                        remaining = total_estimated - elapsed
                        eta_text = f"예상 남은 시간: {remaining / 60:.1f}분" if remaining > 60 else f"예상 남은 시간: {remaining:.0f}초"
                        self._safe_gui_update(lambda eta=eta_text: self.eta_label.configure(text=eta))

                try:
                    result = processor.process_video(
                        video_path=video_path,
                        output_base_dir=self.output_dir_var.get(),
                        progress_callback=progress_callback
                    )

                    if result.errors:
                        self.error_manager.log(LogLevel.WARNING, f"{video_name} 처리 완료 (오류 있음)")
                    else:
                        self.error_manager.log(LogLevel.INFO, f"{video_name} 처리 완료")

                    completed_videos += 1

                except Exception as e:
                    self.error_manager.log(LogLevel.ERROR, f"{video_name} 처리 실패", e)

                overall_progress = ((i + 1) / total_videos) * 100
                self._safe_gui_update(lambda p=overall_progress: self.progress_var.set(p))

            if self.processing:
                self._safe_gui_update(lambda: self.progress_var.set(100))
                self._safe_gui_update(lambda: self.status_label.configure(
                    text=f"✅ 처리 완료! ({completed_videos}/{total_videos})"))
                self._safe_gui_update(lambda: self.current_file_label.configure(text=""))
                self._safe_gui_update(lambda: self.eta_label.configure(text=""))

                if completed_videos > 0:
                    self._safe_gui_update(lambda: messagebox.showinfo(
                        "완료", f"{completed_videos}/{total_videos}개 비디오 처리가 완료되었습니다."))
            else:
                self._safe_gui_update(lambda: self.status_label.configure(text="⏹️ 처리 중지됨"))
                self._safe_gui_update(lambda: self.current_file_label.configure(text=""))
                self._safe_gui_update(lambda: self.eta_label.configure(text=""))

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "비디오 처리 중 오류 발생", e)
            self._safe_gui_update(lambda: messagebox.showerror(
                "오류", f"처리 중 오류가 발생했습니다:\n{str(e)}"))

        finally:
            self.processing = False
            self._safe_gui_update(lambda: self.process_button.configure(state="normal"))
            self._safe_gui_update(lambda: self.stop_button.configure(state="disabled"))

    def on_closing(self):
        """프로그램 종료 처리"""
        try:
            if self.processing:
                if messagebox.askokcancel("종료", "처리가 진행 중입니다. 정말 종료하시겠습니까?"):
                    self.processing = False
                    if self.processing_thread and self.processing_thread.is_alive():
                        self.error_manager.log(LogLevel.INFO, "처리 중지 중...")
                        self.processing_thread.join(timeout=3)
                else:
                    return

            # 설정 저장
            try:
                self.save_config()
            except Exception:
                pass

            self.path_manager.cleanup_all()
            self.root.destroy()

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, "종료 처리 실패", e)
            self.root.destroy()


# ============================================================================
# 오디오 매니저
# ============================================================================

class AudioManager:
    """음성 처리 관리"""

    def __init__(self):
        self.path_manager = PathManager()
        self.error_manager = ErrorManager()
        self.model = None
        self.model_size = "base"

    def extract_audio(self, video_path: Union[str, Path],
                      output_dir: Union[str, Path]) -> Optional[Path]:
        """오디오 추출 - MoviePy 우선, ffmpeg 백업"""
        self.error_manager.log(LogLevel.INFO, f"오디오 추출 시작 - MoviePy 사용 가능: {MOVIEPY_AVAILABLE}")

        # MoviePy 먼저 시도
        if MOVIEPY_AVAILABLE:
            self.error_manager.log(LogLevel.INFO, "MoviePy를 사용하여 오디오 추출 시도...")
            result = self._extract_audio_moviepy(video_path, output_dir)
            if result:
                return result
            self.error_manager.log(LogLevel.WARNING, "MoviePy 오디오 추출 실패, ffmpeg 시도...")
        else:
            self.error_manager.log(LogLevel.WARNING, "MoviePy를 사용할 수 없습니다. ffmpeg로 진행...")

        # ffmpeg 백업
        return self._extract_audio_ffmpeg(video_path, output_dir)

    def _extract_audio_moviepy(self, video_path: Union[str, Path],
                               output_dir: Union[str, Path]) -> Optional[Path]:
        """MoviePy를 사용한 오디오 추출"""
        try:
            video_path = self.path_manager.normalize_path(video_path)
            output_dir = self.path_manager.normalize_path(output_dir)

            video_name = self.path_manager.safe_filename(video_path.stem)
            audio_path = output_dir / f"{video_name}_audio.wav"

            self.error_manager.log(LogLevel.INFO, "MoviePy로 오디오 추출 중...")

            try:
                from moviepy.editor import VideoFileClip
            except ImportError as e:
                self.error_manager.log(LogLevel.ERROR, f"MoviePy import 실패: {e}")
                return None

            # ffmpeg 사용 가능성 체크
            try:
                import imageio
                # ffmpeg 플러그인 확인
                try:
                    imageio.plugins.ffmpeg.get_exe()
                except Exception:
                    self.error_manager.log(LogLevel.WARNING, "imageio-ffmpeg 플러그인이 없습니다. 자동 다운로드 시도...")
                    try:
                        imageio.plugins.ffmpeg.download()
                        self.error_manager.log(LogLevel.INFO, "ffmpeg 플러그인 다운로드 완료")
                    except Exception as download_err:
                        self.error_manager.log(LogLevel.ERROR, f"ffmpeg 플러그인 다운로드 실패: {download_err}")
                        return None
            except ImportError:
                self.error_manager.log(LogLevel.WARNING, "imageio가 없습니다. pip install imageio-ffmpeg 권장")

            with VideoFileClip(str(video_path)) as video:
                if video.audio is None:
                    self.error_manager.log(LogLevel.WARNING, "비디오에 오디오 트랙이 없습니다.")
                    return None

                with video.audio as audio:
                    # 임시 파일 경로 생성 (한글 경로 문제 방지)
                    temp_audio = self.path_manager.create_temp_file(suffix='.wav')

                    audio.write_audiofile(
                        str(audio_path),
                        logger=None,
                        verbose=False,
                        temp_audiofile=str(temp_audio) if temp_audio else None
                    )

            if audio_path.exists() and audio_path.stat().st_size > 0:
                self.error_manager.log(LogLevel.INFO, f"MoviePy 오디오 추출 완료: {audio_path}")
                return audio_path
            else:
                self.error_manager.log(LogLevel.ERROR, "MoviePy로 오디오 파일이 생성되지 않았습니다.")
                return None

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"MoviePy 오디오 추출 실패: {video_path}", e)
            # 구체적인 오류 유형별 안내
            error_str = str(e).lower()
            if 'ffmpeg' in error_str:
                self.error_manager.log(LogLevel.INFO, "💡 해결 방법: pip install imageio-ffmpeg")
            elif 'codec' in error_str:
                self.error_manager.log(LogLevel.INFO, "💡 코덱 문제일 수 있습니다. 다른 비디오로 테스트해보세요.")
            elif 'permission' in error_str:
                self.error_manager.log(LogLevel.INFO, "💡 파일 권한 문제일 수 있습니다.")

            return None

    def _extract_audio_ffmpeg(self, video_path: Union[str, Path],
                              output_dir: Union[str, Path]) -> Optional[Path]:
        """ffmpeg를 사용한 오디오 추출 - 한글 경로 지원"""
        try:
            video_path = self.path_manager.normalize_path(video_path)
            output_dir = self.path_manager.normalize_path(output_dir)

            video_name = self.path_manager.safe_filename(video_path.stem)
            audio_path = output_dir / f"{video_name}_audio.wav"

            # 한글 경로 문제 해결을 위해 임시 파일 사용
            temp_video = None
            temp_audio = self.path_manager.create_temp_file(suffix='.wav')

            if not temp_audio:
                return None

            # 비디오 파일 경로에 한글이 있는지 확인
            try:
                # OpenCV로 테스트
                with self.video_manager.safe_video_capture(video_path) as cap:
                    input_path = str(video_path)
            except:
                # 한글 경로 문제일 가능성 - 임시 복사
                temp_video = self.path_manager.create_temp_file(suffix=video_path.suffix)
                if temp_video and self.path_manager.copy_with_korean_support(video_path, temp_video):
                    input_path = str(temp_video)
                else:
                    return None

            cmd = [
                'ffmpeg', '-i', input_path,
                '-vn', '-acodec', 'pcm_s16le',
                '-ar', '16000', '-ac', '1', '-y',
                str(temp_audio)
            ]

            self.error_manager.log(LogLevel.INFO, "ffmpeg를 사용하여 오디오 추출 중...")

            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=300,
                check=False,
                encoding='utf-8',
                errors='ignore'
            )

            if result.returncode == 0 and temp_audio.exists() and temp_audio.stat().st_size > 0:
                # 최종 위치로 복사
                if self.path_manager.copy_with_korean_support(temp_audio, audio_path):
                    self.error_manager.log(LogLevel.INFO, f"ffmpeg 오디오 추출 완료: {audio_path}")
                    return audio_path

            if result.stderr:
                self.error_manager.log(LogLevel.ERROR, f"ffmpeg 오디오 추출 실패: {result.stderr[:500]}")

            return None

        except subprocess.TimeoutExpired:
            self.error_manager.log(LogLevel.ERROR, "오디오 추출 시간 초과")
            return None
        except FileNotFoundError:
            self.error_manager.log(LogLevel.ERROR, "ffmpeg가 설치되지 않았거나 PATH에 없습니다.")
            return None
        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"오디오 추출 실패: {video_path}", e)
            return None

    def load_whisper_model(self, model_size: str = "base") -> bool:
        """Whisper 모델 로딩"""
        if not WHISPER_AVAILABLE:
            self.error_manager.log(LogLevel.ERROR, "openai-whisper가 설치되지 않았습니다.")
            return False

        if self.model is not None and self.model_size == model_size:
            return True

        try:
            self.error_manager.log(LogLevel.INFO, f"Whisper 모델 로딩: {model_size}")
            self.model = whisper.load_model(model_size)
            self.model_size = model_size
            self.error_manager.log(LogLevel.INFO, f"Whisper 모델 로딩 성공: {model_size}")
            return True

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"Whisper 모델 로딩 실패: {model_size}", e)
            return False

    def transcribe_audio(self, audio_path: Union[str, Path]) -> Optional[Dict]:
        """오디오 텍스트 변환 - 한글 지원"""
        if not WHISPER_AVAILABLE:
            self.error_manager.log(LogLevel.ERROR, "openai-whisper가 설치되지 않았습니다.")
            return None

        try:
            audio_path = self.path_manager.normalize_path(audio_path)

            if not audio_path.exists():
                self.error_manager.log(LogLevel.ERROR, f"오디오 파일이 없습니다: {audio_path}")
                return None

            if audio_path.stat().st_size == 0:
                self.error_manager.log(LogLevel.ERROR, "오디오 파일이 비어있습니다.")
                return None

            if self.model is None:
                if not self.load_whisper_model(self.model_size):
                    return None

            # 한글 경로 문제 해결
            temp_audio = None
            try:
                # 경로에 한글이 있는지 확인
                str(audio_path).encode('ascii')
                input_path = str(audio_path)
            except UnicodeEncodeError:
                # 한글이 있음 - 임시 파일로 복사
                temp_audio = self.path_manager.create_temp_file(suffix='.wav')
                if temp_audio and self.path_manager.copy_with_korean_support(audio_path, temp_audio):
                    input_path = str(temp_audio)
                else:
                    return None

            self.error_manager.log(LogLevel.INFO, "음성 인식 시작...")
            result = self.model.transcribe(
                input_path,
                language="ko",
                verbose=False,
                word_timestamps=True,
                fp16=False
            )

            if not result or 'segments' not in result:
                self.error_manager.log(LogLevel.WARNING, "음성 인식 결과가 비어있습니다.")
                return None

            self.error_manager.log(LogLevel.INFO,
                                   f"음성 인식 완료: {len(result.get('segments', []))}개 세그먼트")
            return result

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"음성 인식 실패: {audio_path}", e)
            return None

    def create_srt(self, transcription: Dict, output_path: Union[str, Path]) -> bool:
        """SRT 자막 파일 생성"""
        if not transcription or 'segments' not in transcription:
            self.error_manager.log(LogLevel.ERROR, "유효하지 않은 전사 데이터")
            return False

        try:
            output_path = self.path_manager.normalize_path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            with open(output_path, 'w', encoding='utf-8') as f:
                segments = transcription.get('segments', [])

                if not segments:
                    self.error_manager.log(LogLevel.WARNING, "전사 세그먼트가 비어있습니다.")
                    return False

                srt_index = 1
                for segment in segments:
                    try:
                        text = segment.get('text', '').strip()
                        if not text:
                            continue

                        start_time = self._seconds_to_srt_time(segment.get('start', 0))
                        end_time = self._seconds_to_srt_time(segment.get('end', 0))

                        # 텍스트 줄 나누기 (가독성 향상)
                        lines = self._split_text_for_srt(text)

                        f.write(f"{srt_index}\n")
                        f.write(f"{start_time} --> {end_time}\n")
                        f.write('\n'.join(lines))
                        f.write("\n\n")

                        srt_index += 1

                    except (KeyError, TypeError) as e:
                        self.error_manager.log(LogLevel.WARNING, f"세그먼트 처리 실패", e)
                        continue

            self.error_manager.log(LogLevel.INFO, f"SRT 파일 생성 완료: {output_path} ({srt_index - 1}개 자막)")
            return True

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"SRT 파일 생성 실패: {output_path}", e)
            return False

    def _split_text_for_srt(self, text: str, max_chars_per_line: int = 40) -> List[str]:
        """SRT용 텍스트 분할"""
        words = text.split()
        lines = []
        current_line = ""

        for word in words:
            if len(current_line) + len(word) + 1 <= max_chars_per_line:
                current_line += (" " + word) if current_line else word
            else:
                if current_line:
                    lines.append(current_line)
                current_line = word

        if current_line:
            lines.append(current_line)

        # 최대 2줄로 제한
        if len(lines) > 2:
            mid_point = len(text) // 2
            lines = [text[:mid_point].strip(), text[mid_point:].strip()]

        return lines

    def _seconds_to_srt_time(self, seconds: float) -> str:
        """초를 SRT 시간 형식으로 변환"""
        try:
            seconds = max(0, seconds)
            td = timedelta(seconds=seconds)
            hours = td.seconds // 3600
            minutes = (td.seconds % 3600) // 60
            secs = td.seconds % 60
            milliseconds = td.microseconds // 1000
            return f"{hours:02d}:{minutes:02d}:{secs:02d},{milliseconds:03d}"
        except Exception:
            return "00:00:00,000"


# ============================================================================
# 문서 매니저
# ============================================================================

class DocumentManager:
    """문서 생성 관리"""

    def __init__(self):
        self.path_manager = PathManager()
        self.error_manager = ErrorManager()
        # SpellChecker 초기화 (사용 가능한 경우)
        if SPELL_CHECKER_AVAILABLE:
            self.spell_checker = SpellChecker(use_ai_model=True)
        else:
            self.spell_checker = None

    def create_word_document(self, transcription: Optional[Dict], frames: List[str],
                             output_path: Union[str, Path], video_name: str) -> bool:
        """Word 문서 생성 - 한글 지원 (transcription 선택적)"""
        if not DOCX_AVAILABLE:
            self.error_manager.log(LogLevel.ERROR, "python-docx가 설치되지 않았습니다.")
            return False

        try:
            output_path = self.path_manager.normalize_path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            doc = Document()

            # 문서 스타일 설정
            self._setup_document_styles(doc)

            # 문서 구조 생성
            self._add_title_section(doc, video_name)
            self._add_document_info_section(doc, frames, transcription)

            if transcription and transcription.get('segments'):
                self._add_full_transcript_section(doc, transcription)
                self._add_timestamped_transcript_section(doc, transcription)

            if frames:
                self._add_frames_gallery_section(doc, frames)
            elif not transcription:
                # 프레임도 없고 음성인식도 없으면 경고
                self.error_manager.log(LogLevel.WARNING, "생성할 콘텐츠가 없습니다 (프레임 또는 음성인식 데이터 필요)")

            # 문서 저장
            doc.save(str(output_path))
            self.error_manager.log(LogLevel.INFO, f"Word 문서 생성 완료: {output_path}")
            return True

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"Word 문서 생성 실패: {output_path}", e)
            return False

    def _setup_document_styles(self, doc):
        """문서 스타일 설정"""
        try:
            style = doc.styles['Normal']
            style.font.name = '맑은 고딕'
            style.font.size = DocxPt(11)
            style.paragraph_format.line_spacing = 1.5
        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "문서 스타일 설정 실패", e)

    def _add_title_section(self, doc, video_name: str):
        """제목 섹션 추가"""
        try:
            title = doc.add_heading(f'{video_name}', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            subtitle = doc.add_heading('비디오 분석 보고서', level=1)
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

            doc.add_paragraph()
        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "제목 섹션 추가 실패", e)

    def _add_document_info_section(self, doc, frames: List[str], transcription: Optional[Dict]):
        """문서 정보 섹션 추가"""
        try:
            info_heading = doc.add_heading('📋 문서 정보', level=2)

            info_table = doc.add_table(rows=0, cols=2)
            info_table.style = 'Light List Accent 1'

            # 생성일
            row = info_table.add_row()
            row.cells[0].text = '생성일'
            row.cells[1].text = datetime.now().strftime('%Y년 %m월 %d일 %H:%M')

            # 총 프레임
            row = info_table.add_row()
            row.cells[0].text = '추출된 프레임'
            row.cells[1].text = f'{len(frames)}개'

            # 영상 정보
            if transcription and transcription.get('segments'):
                duration = transcription.get('segments', [{}])[-1].get('end', 0)
                row = info_table.add_row()
                row.cells[0].text = '영상 길이'
                row.cells[1].text = self._format_duration(duration)

                row = info_table.add_row()
                row.cells[0].text = '음성 세그먼트'
                row.cells[1].text = f"{len(transcription.get('segments', []))}개"

            doc.add_paragraph()
            doc.add_page_break()

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "문서 정보 섹션 추가 실패", e)

    def _add_full_transcript_section(self, doc, transcription: Dict):
        """전체 자막 섹션 추가"""
        try:
            doc.add_heading('📝 전체 자막 (Full Transcript)', level=1)
            doc.add_paragraph()

            segments = transcription.get('segments', [])
            if not segments:
                doc.add_paragraph("전사된 텍스트가 없습니다.")
                return

            # 문단 단위로 그룹화
            paragraph_groups = self._group_segments_into_paragraphs(segments)

            for i, paragraph_text in enumerate(paragraph_groups):
                if paragraph_text:
                    p = doc.add_paragraph(paragraph_text)
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    p.paragraph_format.first_line_indent = DocxPt(20)
                    p.paragraph_format.space_after = DocxPt(12)

                    if (i + 1) % 5 == 0 and i < len(paragraph_groups) - 1:
                        doc.add_paragraph('─' * 50)
                        doc.add_paragraph()

            doc.add_page_break()

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "전체 자막 섹션 추가 실패", e)

    def _add_timestamped_transcript_section(self, doc, transcription: Dict):
        """시간별 자막 섹션 추가"""
        try:
            doc.add_heading('⏱️ 시간별 자막 (Timestamped Transcript)', level=1)
            doc.add_paragraph()

            segments = transcription.get('segments', [])
            time_groups = self._group_segments_by_time(segments)

            for minute in sorted(time_groups.keys()):
                heading = doc.add_heading(f'{minute}분 ~ {minute + 1}분', level=3)

                table = doc.add_table(rows=1, cols=2)
                table.style = 'Light Grid Accent 1'

                hdr_cells = table.rows[0].cells
                hdr_cells[0].text = '시간'
                hdr_cells[1].text = '내용'

                for cell in hdr_cells:
                    cell.paragraphs[0].runs[0].font.bold = True
                    cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

                for segment in time_groups[minute]:
                    try:
                        start_time = self._format_time(segment.get('start', 0))
                        text = segment.get('text', '').strip()

                        if text:
                            row_cells = table.add_row().cells
                            row_cells[0].text = start_time
                            row_cells[1].text = text

                            row_cells[0].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                            row_cells[1].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.LEFT

                    except Exception as e:
                        self.error_manager.log(LogLevel.WARNING, "세그먼트 처리 실패", e)
                        continue

                doc.add_paragraph()

            doc.add_page_break()

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "시간별 자막 섹션 추가 실패", e)

    def _add_frames_gallery_section(self, doc, frames: List[str]):
        """프레임 갤러리 섹션 추가"""
        try:
            doc.add_heading('🎬 주요 프레임 (Key Frames)', level=1)
            doc.add_paragraph()

            p = doc.add_paragraph(f"총 {len(frames)}개의 주요 프레임이 추출되었습니다.")
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph()

            # 이미지 갤러리 생성 (3열 레이아웃)
            max_frames = min(len(frames), 20)
            cols_per_row = 3

            for i in range(0, max_frames, cols_per_row):
                table = doc.add_table(rows=2, cols=cols_per_row)
                table.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # 이미지 행
                for j in range(cols_per_row):
                    if i + j < max_frames:
                        cell = table.rows[0].cells[j]
                        frame_path = frames[i + j]

                        try:
                            frame_path_obj = Path(frame_path)
                            if frame_path_obj.exists():
                                paragraph = cell.paragraphs[0]
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                run = paragraph.add_run()

                                # 이미지 최적화 및 추가
                                temp_image = self._optimize_image_for_word(frame_path_obj)
                                if temp_image:
                                    run.add_picture(str(temp_image), width=DocxInches(2.0))
                                else:
                                    cell.text = "이미지 로드 실패"
                            else:
                                cell.text = "이미지 없음"

                        except Exception as e:
                            self.error_manager.log(LogLevel.WARNING, f"이미지 추가 실패 {frame_path}", e)
                            cell.text = "이미지 로드 실패"
                    else:
                        table.rows[0].cells[j].text = ""

                # 캡션 행
                for j in range(cols_per_row):
                    if i + j < max_frames:
                        cell = table.rows[1].cells[j]
                        frame_path = frames[i + j]

                        caption = self._extract_time_from_filename(Path(frame_path).name)

                        p = cell.paragraphs[0]
                        p.text = caption
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        if p.runs:
                            p.runs[0].font.size = DocxPt(9)
                            p.runs[0].font.color.rgb = RGBColor(100, 100, 100)
                    else:
                        table.rows[1].cells[j].text = ""

                if i + cols_per_row < max_frames:
                    doc.add_paragraph()

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "프레임 갤러리 섹션 추가 실패", e)

    def _optimize_image_for_word(self, image_path: Path) -> Optional[Path]:
        """Word용 이미지 최적화"""
        try:
            with Image.open(image_path) as img:
                if img.width > 1920 or img.height > 1080:
                    temp_path = self.path_manager.create_temp_file(suffix='.jpg')
                    if temp_path:
                        img.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                        img.save(temp_path, quality=85)
                        return temp_path
                else:
                    return image_path
        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, f"이미지 최적화 실패: {image_path}", e)
            return None

    def _group_segments_into_paragraphs(self, segments: List[Dict],
                                        time_threshold: float = 30,
                                        segment_threshold: int = 5) -> List[str]:
        """세그먼트를 문단으로 그룹화"""
        paragraph_groups = []
        current_group = []
        current_start_time = 0

        for segment in segments:
            if not current_group:
                current_start_time = segment.get('start', 0)

            current_group.append(segment.get('text', '').strip())

            time_diff = segment.get('end', 0) - current_start_time
            if time_diff >= time_threshold or len(current_group) >= segment_threshold:
                if current_group:
                    paragraph_groups.append(' '.join(current_group))
                    current_group = []

        if current_group:
            paragraph_groups.append(' '.join(current_group))

        return paragraph_groups

    def _group_segments_by_time(self, segments: List[Dict]) -> Dict[int, List[Dict]]:
        """세그먼트를 시간대별로 그룹화"""
        time_groups = {}
        for segment in segments:
            start_time = segment.get('start', 0)
            minute = int(start_time // 60)

            if minute not in time_groups:
                time_groups[minute] = []

            time_groups[minute].append(segment)

        return time_groups

    def _extract_time_from_filename(self, filename: str) -> str:
        """파일명에서 시간 정보 추출"""
        try:
            if 'frame_' in filename:
                # video_name_frame_0001s_000123.jpg 형식
                parts = filename.split('_')
                for part in parts:
                    if part.endswith('s'):
                        time_sec = int(part[:-1])
                        return f"{time_sec}초"
            return filename.replace('.jpg', '').replace('_frame_', '_')
        except:
            return filename.replace('.jpg', '')

    def _format_time(self, seconds: float) -> str:
        """초를 시:분:초 형식으로 변환"""
        try:
            seconds = max(0, seconds)
            hours = int(seconds // 3600)
            minutes = int((seconds % 3600) // 60)
            secs = int(seconds % 60)

            if hours > 0:
                return f"{hours:02d}:{minutes:02d}:{secs:02d}"
            else:
                return f"{minutes:02d}:{secs:02d}"
        except Exception:
            return "00:00"

    def _format_duration(self, seconds: float) -> str:
        """초를 읽기 쉬운 형식으로 변환"""
        try:
            seconds = max(0, seconds)
            hours = int(seconds // 3600)
            minutes = int((seconds % 3600) // 60)
            secs = int(seconds % 60)

            parts = []
            if hours > 0:
                parts.append(f"{hours}시간")
            if minutes > 0:
                parts.append(f"{minutes}분")
            if secs > 0 or not parts:
                parts.append(f"{secs}초")

            return ' '.join(parts)
        except Exception:
            return "0초"


# ============================================================================
# PPT 매니저
# ============================================================================

class PPTManager:
    """PPT 생성 관리"""

    def __init__(self):
        self.path_manager = PathManager()
        self.error_manager = ErrorManager()

    def create_ppt(self, frames: List[str], output_path: Union[str, Path],
                   video_name: str, grid_size: Tuple[int, int]) -> bool:
        """PPT 생성 - 한글 지원"""
        if not PPTX_AVAILABLE:
            self.error_manager.log(LogLevel.WARNING, "python-pptx가 설치되지 않아 PPT를 생성할 수 없습니다.")
            return False

        if not frames:
            self.error_manager.log(LogLevel.WARNING, "저장된 프레임이 없어 PPT를 생성할 수 없습니다.")
            return False

        try:
            output_path = self.path_manager.normalize_path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            rows, cols = grid_size
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            # 슬라이드 레이아웃 계산
            images_per_slide = rows * cols
            margin = Inches(0.2)
            spacing = Inches(0.1)
            img_width = (prs.slide_width - 2 * margin - (cols - 1) * spacing) / cols
            img_height = (prs.slide_height - 2 * margin - (rows - 1) * spacing - Inches(0.4)) / rows

            # 타이틀 슬라이드
            self._add_title_slide(prs, video_name, len(frames))

            # 이미지 슬라이드들
            total_images = len(frames)
            total_slides = math.ceil(total_images / images_per_slide)

            self.error_manager.log(LogLevel.INFO, f"PPT 생성 중... (슬라이드: {total_slides}개)")

            for slide_idx in range(total_slides):
                self._add_image_slide(prs, frames, slide_idx, images_per_slide,
                                      video_name, total_slides, rows, cols,
                                      margin, spacing, img_width, img_height)

            # PPT 저장
            prs.save(str(output_path))
            self.error_manager.log(LogLevel.INFO, f"PPT 생성 완료: {output_path}")
            return True

        except Exception as e:
            self.error_manager.log(LogLevel.ERROR, f"PPT 생성 실패: {output_path}", e)
            return False

    def _add_title_slide(self, prs, video_name: str, frame_count: int):
        """타이틀 슬라이드 추가"""
        try:
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = f"{video_name} - 프레임 캡처"
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = (
                    f"총 {frame_count}개 프레임\n"
                    f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
                )
        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, "타이틀 슬라이드 추가 실패", e)

    def _add_image_slide(self, prs, frames: List[str], slide_idx: int,
                         images_per_slide: int, video_name: str, total_slides: int,
                         rows: int, cols: int, margin, spacing, img_width, img_height):
        """이미지 슬라이드 추가"""
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # 슬라이드 제목
            title_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(0.1),
                width=prs.slide_width - Inches(1),
                height=Inches(0.3)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"{video_name} - 슬라이드 {slide_idx + 1}/{total_slides}"
            title_frame.paragraphs[0].font.size = Pt(16)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            start_idx = slide_idx * images_per_slide
            end_idx = min(start_idx + images_per_slide, len(frames))

            for img_idx in range(start_idx, end_idx):
                grid_position = img_idx - start_idx
                row = grid_position // cols
                col = grid_position % cols

                left = margin + col * (img_width + spacing)
                top = margin + Inches(0.4) + row * (img_height + spacing)

                frame_path = Path(frames[img_idx])
                if not frame_path.exists():
                    self.error_manager.log(LogLevel.WARNING, f"프레임 파일이 없습니다: {frame_path}")
                    continue

                try:
                    # 이미지 최적화 및 추가
                    temp_image = self._optimize_image_for_ppt(frame_path)
                    if temp_image:
                        slide.shapes.add_picture(str(temp_image), left, top, img_width, img_height)

                        # 파일명 캡션 추가
                        self._add_image_caption(slide, frame_path, left, top, img_width, img_height)

                except Exception as e:
                    self.error_manager.log(LogLevel.ERROR, f"PPT 이미지 추가 실패 {frame_path}", e)
                    continue

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, f"이미지 슬라이드 추가 실패: 슬라이드 {slide_idx + 1}", e)

    def _optimize_image_for_ppt(self, image_path: Path) -> Optional[Path]:
        """PPT용 이미지 최적화"""
        try:
            with Image.open(image_path) as img:
                if img.width > 1920 or img.height > 1080:
                    temp_path = self.path_manager.create_temp_file(suffix='.jpg')
                    if temp_path:
                        img.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                        img.save(temp_path, quality=85, optimize=True)
                        return temp_path
                else:
                    return image_path
        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, f"이미지 최적화 실패: {image_path}", e)
            return None

    def _add_image_caption(self, slide, frame_path: Path, left, top, img_width, img_height):
        """이미지 캡션 추가"""
        try:
            caption = self._extract_time_from_filename(frame_path.name)

            text_box = slide.shapes.add_textbox(
                left=left,
                top=top + img_height,
                width=img_width,
                height=Inches(0.2)
            )
            text_frame = text_box.text_frame
            text_frame.text = caption
            text_frame.paragraphs[0].font.size = Pt(8)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        except Exception as e:
            self.error_manager.log(LogLevel.WARNING, f"캡션 추가 실패: {frame_path}", e)

    def _extract_time_from_filename(self, filename: str) -> str:
        """파일명에서 시간 정보 추출"""
        try:
            if 'frame_' in filename:
                parts = filename.split('_')
                for part in parts:
                    if part.endswith('s'):
                        time_sec = int(part[:-1])
                        return f"{time_sec}초"
            return filename.replace('.jpg', '').replace('_frame_', '_')
        except:
            return filename.replace('.jpg', '')


# ============================================================================
# 메인 처리 클래스
# ============================================================================

class VideoProcessorMain:
    """메인 비디오 처리 클래스"""

    def __init__(self):
        self.config_manager = ConfigManager()
        self.error_manager = ErrorManager()
        self.path_manager = PathManager()

        # 각 매니저 인스턴스
        self.video_manager = VideoManager()
        self.frame_extraction_manager = FrameExtractionManager()
        self.audio_manager = AudioManager()
        self.document_manager = DocumentManager()
        self.ppt_manager = PPTManager()

        # OCRManager (사용 가능한 경우)
        if OCR_AVAILABLE:
            self.ocr_manager = OCRManager(use_ai_model=True)
        else:
            self.ocr_manager = None

    def process_video(self, video_path: Union[str, Path],
                      output_base_dir: Union[str, Path] = "output",
                      progress_callback=None) -> ProcessingResult:
        """비디오 통합 처리"""

        start_time = time.time()
        video_path = self.path_manager.normalize_path(video_path)

        result = ProcessingResult(
            video_path=str(video_path),
            status=ProcessingStatus.PROCESSING
        )

        try:
            # 입력 검증
            if not video_path.exists():
                error_msg = f"비디오 파일이 존재하지 않습니다: {video_path}"
                self.error_manager.log(LogLevel.ERROR, error_msg)
                result.errors.append(error_msg)
                result.status = ProcessingStatus.ERROR
                return result

            # 출력 디렉토리 설정
            video_name = self.path_manager.safe_filename(video_path.stem)
            output_dir = self.path_manager.normalize_path(output_base_dir) / video_name

            if not self.path_manager.ensure_directory(output_dir):
                error_msg = f"출력 디렉토리 생성 실패: {output_dir}"
                self.error_manager.log(LogLevel.ERROR, error_msg)
                result.errors.append(error_msg)
                result.status = ProcessingStatus.ERROR
                return result

            result.output_dir = str(output_dir)
            config = self.config_manager.get_config()

            self.error_manager.log(LogLevel.INFO, "=" * 60)
            self.error_manager.log(LogLevel.INFO, f"비디오 처리 시작: {video_path.name}")
            self.error_manager.log(LogLevel.INFO, "=" * 60)

            # 1. 프레임 추출
            if config.extract_frames:
                self.error_manager.log(LogLevel.INFO, "📷 프레임 추출 중...")
                try:
                    frames, frame_count = self.frame_extraction_manager.extract_frames(
                        video_path, output_dir, progress_callback
                    )
                    result.frames = frames
                    result.frame_count = frame_count

                    if frame_count == 0:
                        result.warnings.append("프레임이 추출되지 않았습니다")

                except Exception as e:
                    error_msg = f"프레임 추출 실패: {e}"
                    self.error_manager.log(LogLevel.ERROR, error_msg, e)
                    result.errors.append(error_msg)

            # 2. PPT 생성
            if config.create_ppt and result.frames:
                self.error_manager.log(LogLevel.INFO, "📊 PPT 생성 중...")
                try:
                    ppt_path = output_dir / f"{video_name}_frames.pptx"
                    if self.ppt_manager.create_ppt(
                            result.frames, ppt_path, video_name,
                            (config.grid_rows, config.grid_cols)
                    ):
                        result.ppt_path = str(ppt_path)
                    else:
                        result.warnings.append("PPT 생성 실패")

                except Exception as e:
                    error_msg = f"PPT 생성 실패: {e}"
                    self.error_manager.log(LogLevel.ERROR, error_msg, e)
                    result.errors.append(error_msg)

            # 3. 음성 처리
            if config.extract_audio:
                self.error_manager.log(LogLevel.INFO, "🎵 오디오 추출 중...")
                try:
                    # Whisper 모델 설정
                    self.audio_manager.model_size = config.whisper_model

                    audio_path = self.audio_manager.extract_audio(video_path, output_dir)
                    result.audio_path = str(audio_path) if audio_path else None

                    if not audio_path:
                        result.warnings.append("오디오 추출 실패")
                    elif config.create_srt or config.create_word:
                        self.error_manager.log(LogLevel.INFO, "🎙️ 음성 인식 중...")
                        transcription = self.audio_manager.transcribe_audio(audio_path)
                        result.transcription = transcription

                        if not transcription:
                            result.warnings.append("음성 인식 실패")
                        else:
                            # 4. SRT 생성
                            if config.create_srt:
                                self.error_manager.log(LogLevel.INFO, "📝 SRT 자막 생성 중...")
                                try:
                                    srt_path = output_dir / f"{video_name}.srt"
                                    if self.audio_manager.create_srt(transcription, srt_path):
                                        result.srt_path = str(srt_path)
                                    else:
                                        result.warnings.append("SRT 생성 실패")

                                except Exception as e:
                                    error_msg = f"SRT 생성 실패: {e}"
                                    self.error_manager.log(LogLevel.ERROR, error_msg, e)
                                    result.errors.append(error_msg)

                except Exception as e:
                    error_msg = f"음성 처리 실패: {e}"
                    self.error_manager.log(LogLevel.ERROR, error_msg, e)
                    result.errors.append(error_msg)

            # 4. Word 문서 생성 (오디오와 독립적으로 실행)
            if config.create_word:
                self.error_manager.log(LogLevel.INFO, "📄 Word 문서 생성 중...")
                try:
                    word_path = output_dir / f"{video_name}_report.docx"
                    # transcription이 없으면 None으로 전달 (프레임만으로 문서 생성)
                    transcription_data = result.transcription if hasattr(result, 'transcription') else None
                    if self.document_manager.create_word_document(
                            transcription_data, result.frames, word_path, video_name
                    ):
                        result.word_path = str(word_path)
                    else:
                        result.warnings.append("Word 문서 생성 실패")

                except Exception as e:
                    error_msg = f"Word 문서 생성 실패: {e}"
                    self.error_manager.log(LogLevel.ERROR, error_msg, e)
                    result.errors.append(error_msg)

            # 6. 프레임 OCR 검수 (선택적 - 설정으로 활성화)
            if config.ocr_spell_check and config.extract_frames and result.frames and self.ocr_manager and OCR_AVAILABLE:
                self.error_manager.log(LogLevel.INFO, "🔍 프레임 OCR 검수 중...")
                try:
                    # 프레임 폴더에서 OCR 검수 실행
                    frames_dir = Path(output_dir) / "frames"
                    if frames_dir.exists():
                        ocr_output_dir = Path(output_dir) / "frames_ocr_checked"
                        ocr_output_dir.mkdir(exist_ok=True)

                        def ocr_progress(current, total, path):
                            if progress_callback:
                                progress_callback(f"OCR 검수: {current}/{total}")

                        ocr_result = self.ocr_manager.process_folder(
                            frames_dir,
                            output_dir=ocr_output_dir,
                            comparison_mode=True,
                            file_pattern="*.jpg",
                            callback=ocr_progress
                        )

                        if ocr_result['success']:
                            self.error_manager.log(LogLevel.INFO,
                                f"✅ OCR 검수 완료: {ocr_result['processed']}개 처리됨")
                            result.warnings.append(
                                f"OCR 검수: {ocr_result['processed']}개 프레임 처리, "
                                f"{ocr_result['failed']}개 실패"
                            )
                        else:
                            result.warnings.append("OCR 검수 실패")

                except Exception as e:
                    error_msg = f"OCR 검수 중 오류: {e}"
                    self.error_manager.log(LogLevel.WARNING, error_msg)
                    result.warnings.append(error_msg)

            # 처리 완료
            result.processing_time = time.time() - start_time

            if result.errors:
                result.status = ProcessingStatus.ERROR
                self.error_manager.log(LogLevel.WARNING,
                                       f"⚠️ 비디오 처리 완료 (오류 {len(result.errors)}개, "
                                       f"경고 {len(result.warnings)}개)")
            else:
                result.status = ProcessingStatus.COMPLETED
                self.error_manager.log(LogLevel.INFO, "✅ 비디오 처리 완료!")

            self.error_manager.log(LogLevel.INFO, f"처리 시간: {result.processing_time:.1f}초")
            self.error_manager.log(LogLevel.INFO, "=" * 60)

            return result

        except Exception as e:
            result.processing_time = time.time() - start_time
            error_msg = f"비디오 처리 중 예상치 못한 오류: {e}"
            self.error_manager.log(LogLevel.ERROR, error_msg, e)
            result.errors.append(error_msg)
            result.status = ProcessingStatus.ERROR
            return result


# ============================================================================
# CLI 및 메인 함수
# ============================================================================

def main_cli():
    """CLI 모드 메인 함수"""
    parser = argparse.ArgumentParser(
        description='올인원 비디오 처리기 v4.0 - 리팩토링 버전',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
사용 예시:
  %(prog)s video1.mp4 video2.mp4 --output results
  %(prog)s *.mp4 --threshold 0.9 --no-audio
  %(prog)s video.mp4 --sampling-interval 2.0 --max-frames 100
  %(prog)s video.mp4 --gui
        """
    )

    parser.add_argument('video_paths', nargs='*', help='처리할 비디오 파일')
    parser.add_argument('--output', '-o', default='output', help='출력 디렉토리')
    parser.add_argument('--threshold', '-t', type=float, default=0.95,
                        help='유사도 임계값 0.0-1.0 (기본: 0.95)')
    parser.add_argument('--sampling-interval', type=float, default=1.0,
                        help='샘플링 간격 (초, 기본: 1.0)')
    parser.add_argument('--min-interval', type=float, default=0.5,
                        help='최소 프레임 간격 (초, 기본: 0.5)')
    parser.add_argument('--max-frames', type=int, default=200,
                        help='비디오당 최대 프레임 수 (기본: 200)')
    parser.add_argument('--no-frames', action='store_true', help='프레임 추출 안 함')
    parser.add_argument('--no-ppt', action='store_true', help='PPT 생성 안 함')
    parser.add_argument('--no-audio', action='store_true', help='음성 추출 안 함')
    parser.add_argument('--no-srt', action='store_true', help='SRT 생성 안 함')
    parser.add_argument('--no-word', action='store_true', help='Word 생성 안 함')
    parser.add_argument('--no-adaptive', action='store_true', help='적응형 임계값 사용 안 함')
    parser.add_argument('--grid-rows', type=int, default=3, choices=range(1, 11),
                        help='PPT 그리드 행 수')
    parser.add_argument('--grid-cols', type=int, default=3, choices=range(1, 11),
                        help='PPT 그리드 열 수')
    parser.add_argument('--whisper-model', default='base',
                        choices=['tiny', 'base', 'small', 'medium', 'large'],
                        help='Whisper 모델 크기')
    parser.add_argument('--gui', action='store_true', help='GUI 모드 실행')
    parser.add_argument('--version', action='version', version='%(prog)s 4.0')

    args = parser.parse_args()

    if args.gui:
        if GUI_AVAILABLE:
            ErrorManager().log(LogLevel.INFO, "GUI 모드로 시작합니다.")
            root = tk.Tk()
            app = VideoProcessorGUI(root)
            try:
                root.mainloop()
            except KeyboardInterrupt:
                ErrorManager().log(LogLevel.INFO, "사용자에 의해 중단되었습니다.")
            finally:
                PathManager().cleanup_all()
        else:
            print("GUI를 사용할 수 없습니다. tkinter가 설치되지 않았습니다.")
            sys.exit(1)
        return

    if not args.video_paths:
        print("오류: 처리할 비디오 파일을 지정하세요.")
        print("GUI 모드: python script.py --gui")
        sys.exit(1)

    # 입력 검증
    if not (0.0 <= args.threshold <= 1.0):
        ErrorManager().log(LogLevel.ERROR, f"유사도 임계값은 0과 1 사이여야 합니다: {args.threshold}")
        sys.exit(1)

    if args.sampling_interval < args.min_interval:
        ErrorManager().log(LogLevel.ERROR, "샘플링 간격은 최소 간격보다 커야 합니다")
        sys.exit(1)

    missing_files = []
    for video_path in args.video_paths:
        if not Path(video_path).exists():
            missing_files.append(video_path)

    if missing_files:
        ErrorManager().log(LogLevel.ERROR, f"존재하지 않는 파일: {', '.join(missing_files)}")
        sys.exit(1)

    try:
        # 설정 적용
        config_manager = ConfigManager()
        config_manager.update_config(
            similarity_threshold=args.threshold,
            adaptive_threshold=not args.no_adaptive,
            extract_frames=not args.no_frames,
            create_ppt=not args.no_ppt,
            extract_audio=not args.no_audio,
            create_srt=not args.no_srt,
            create_word=not args.no_word,
            grid_rows=args.grid_rows, 
            grid_cols=args.grid_cols,
            whisper_model=args.whisper_model,
            frame_sampling_interval=args.sampling_interval,
            min_frame_interval=args.min_interval,
            max_frames_per_video=args.max_frames
        )

        processor = VideoProcessorMain()

        ErrorManager().log(LogLevel.INFO, f"처리할 비디오: {len(args.video_paths)}개")
        ErrorManager().log(LogLevel.INFO, f"출력 디렉토리: {args.output}")

        successful_count = 0

        for i, video_path in enumerate(args.video_paths, 1):
            ErrorManager().log(LogLevel.INFO,
                               f"\n[{i}/{len(args.video_paths)}] 처리 중: {Path(video_path).name}")

            try:
                result = processor.process_video(
                    video_path=video_path,
                    output_base_dir=args.output
                )

                print(f"\n📋 처리 결과: {Path(video_path).name}")
                print(f"   📁 출력: {result.output_dir}")
                print(f"   ⏱️ 처리 시간: {result.processing_time:.1f}초")

                if result.frame_count > 0:
                    print(f"   📷 프레임: {result.frame_count}개")
                if result.ppt_path:
                    print(f"   📊 PPT: {Path(result.ppt_path).name}")
                if result.srt_path:
                    print(f"   📝 자막: {Path(result.srt_path).name}")
                if result.word_path:
                    print(f"   📄 문서: {Path(result.word_path).name}")

                if result.warnings:
                    print(f"   ⚠️ 경고: {len(result.warnings)}개")
                    for warning in result.warnings[:3]:
                        print(f"      - {warning}")

                if result.errors:
                    print(f"   ❌ 오류: {len(result.errors)}개")
                    for error in result.errors[:3]:
                        print(f"      - {error}")

                if result.status == ProcessingStatus.COMPLETED:
                    successful_count += 1

            except KeyboardInterrupt:
                ErrorManager().log(LogLevel.INFO, "\n사용자에 의해 중단되었습니다.")
                break
            except Exception as e:
                ErrorManager().log(LogLevel.ERROR, f"처리 실패: {video_path}", e)

        print(f"\n🎯 전체 결과: {successful_count}/{len(args.video_paths)}개 성공")

    except KeyboardInterrupt:
        ErrorManager().log(LogLevel.INFO, "사용자에 의해 중단되었습니다.")
    except Exception as e:
        ErrorManager().log(LogLevel.ERROR, "프로그램 실행 중 오류", e)
        sys.exit(1)
    finally:
        PathManager().cleanup_all()


def main():
    """메인 진입점"""
    try:
        if len(sys.argv) == 1 and GUI_AVAILABLE:
            ErrorManager().log(LogLevel.INFO, "GUI 모드로 시작합니다.")
            root = tk.Tk()
            app = VideoProcessorGUI(root)
            try:
                root.mainloop()
            except KeyboardInterrupt:
                ErrorManager().log(LogLevel.INFO, "사용자에 의해 중단되었습니다.")
            finally:
                PathManager().cleanup_all()
        else:
            main_cli()

    except Exception as e:
        ErrorManager().log(LogLevel.ERROR, "프로그램 시작 실패", e)
        sys.exit(1)


if __name__ == "__main__":
    main()
